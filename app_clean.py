  # app_clean.py — IC-LicAI Expert Console (Structural + FRAND + LIP Assistant)
# Adds: DOCX/PPTX extraction, weighted IC signal engine, interpreted narrative,
# radar dashboard, CSV semantic extraction, robust company-context auto-split,
# IAS 38 Structural Capital emphasis, FRAND-aware licensing templates,
# Seven Stakeholder / ESG narrative, LIP Console, and LIP Assistant (beta).

from __future__ import annotations
import io, os, tempfile, re, csv
from pathlib import Path
from typing import Dict, Any, List, Tuple, Optional

import streamlit as st
import plotly.graph_objects as go  # for radar charts

# -------------------- PROJECT LOGOS -----------------
IMPACT3T_LOGO_PATH = "demo_assets/impact3t_logo.png"
EU_FLAG_PATH = "demo_assets/eu_flag.png"

# -------------------- MODE / AUTH --------------------
# APP_MODE is controlled by an environment variable on Streamlit:
#   "PRIVATE" = internal app (passphrase + server save)
#   "PUBLIC"  = demo app (no passphrase, no server save)
APP_MODE = os.environ.get("ICLICAI_MODE", "PRIVATE").upper()

PUBLIC_MODE: bool = (APP_MODE == "PUBLIC")
REQUIRE_PASS: bool = (APP_MODE == "PRIVATE")

# ---------------- DOCX/PPTX optional ----------------
HAVE_DOCX = False
HAVE_PPTX = False
try:
    from docx import Document  # type: ignore
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False

try:
    from pptx import Presentation  # type: ignore
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

# ------------------ THEME ----------------------------
# IMPAC3T-IP inspired palette (no yellow / gold)
PRIMARY_NAVY   = "#003B70"  # deep blue, EU-friendly
ACCENT_BLUE    = "#00A7E1"  # bright accent from IMPAC3T-style slides
ACCENT_PURPLE  = "#6E4BAF"  # secondary accent if needed later
BG_MAIN        = "#F5F7FB"  # soft neutral background
SIDEBAR_BG     = "#0F2F56"  # dark blue sidebar
CARD_BG        = "#FFFFFF"  # white cards / blocks
TEXT_MAIN      = "#1F2933"  # dark grey text

st.set_page_config(page_title="IC-LicAI Expert Console", layout="wide")

st.markdown(
    f"""
<style>
  /* Overall app background */
  .stApp {{
    background: {BG_MAIN};
  }}

  /* Main content container */
  .block-container {{
    max-width: 1250px;
    padding-top: 1.2rem;
    padding-bottom: 2rem;
  }}

  /* Title bar at the top */
  .ic-title-bar {{
    background: {PRIMARY_NAVY};
    color: #FFFFFF;
    font-weight: 800;
    font-size: 32px;
    padding: 16px 22px;
    border-radius: 10px;
    letter-spacing: 0.2px;
    margin: 10px 0 20px 0;
    box-shadow: 0 2px 6px rgba(0,0,0,0.08);
  }}

  /* Primary buttons */
  .stButton>button {{
    background: {ACCENT_BLUE} !important;
    color: #FFFFFF !important;
    border-radius: 8px !important;
    border: 0 !important;
    padding: 0.55rem 1rem !important;
    font-weight: 700 !important;
  }}

  .stButton>button:hover {{
    background: {PRIMARY_NAVY} !important;
  }}

  /* Sidebar styling */
  section[data-testid="stSidebar"] {{
    background: {SIDEBAR_BG};
  }}

  section[data-testid="stSidebar"] h1,
  section[data-testid="stSidebar"] h2,
  section[data-testid="stSidebar"] h3,
  section[data-testid="stSidebar"] p,
  section[data-testid="stSidebar"] label,
  section[data-testid="stSidebar"] span {{
    color: #E7F0FF !important;
  }}

  .stRadio div[role="radiogroup"] label {{
    color: #E7F0FF !important;
  }}

  /* Make cards / expanders look clean on light background */
  .stExpander, .stTabs [data-baseweb="tab"] {{
    background: {CARD_BG};
  }}
</style>
""",
    unsafe_allow_html=True,
)

st.markdown(
    '<div class="ic-title-bar">IC-LicAI Expert Console</div>',
    unsafe_allow_html=True,
)
st.caption("INTERNAL VERSION — FOR REAL EVIDENCE (PASS-PHRASE PROTECTED)")
st.markdown(
    '<div class="ic-title-bar">IC-LicAI Expert Console</div>',
    unsafe_allow_html=True,
)
st.caption("INTERNAL VERSION — FOR REAL EVIDENCE (PASS-PHRASE PROTECTED)")

st.markdown("## Belgian Tax Planning Note")

st.info(
    """
    Belgium has announced a new capital gains tax framework for certain financial assets from 1 January 2026.
    For demo purposes, this tool highlights the value of obtaining a robust, evidence-based valuation position
    before a future transaction, while following a data-minimisation approach and not retaining client files.
    """
)

st.markdown(
    """
    **Illustrative planning message**
    
    A stronger, evidence-based valuation record can help support a more defensible starting position when
    assessing future share value movements, exits, restructurings, or shareholder tax exposure.

    In practice, this means:
    - establishing a clear and supportable company value before 2026;
    - strengthening documented Structural Capital (contracts, datasets, software, methods) that can be recognised and defended;
    - using licensing strategies to generate value from these assets before any exit event.

    Licensing can play a critical role by:
    - demonstrating real market value through revenue or access agreements;
    - creating multiple value streams (commercial, partner, and community access);
    - supporting a higher and more defensible valuation position ahead of sale, investment, or restructuring.

    This note is informational only. It is not tax or legal advice.
    """
)

st.markdown(
    """
    **Illustrative planning message**
    
    A stronger, evidence-based valuation record can help support a more defensible starting position when
    assessing future share value movements, exits, restructurings, or shareholder tax exposure.

    This note is informational only. It is not tax or legal advice.
    """
)
# ------------------ AUTH GATE ------------------------
def _auth_gate() -> None:
    if not REQUIRE_PASS:
        return
    secret = st.secrets.get("APP_KEY", None) or os.environ.get("APP_KEY", None)
    if not secret:
        with st.expander("Access control"):
            st.info("Optional passphrase: set st.secrets['APP_KEY'] or env APP_KEY.")
        return
    key = st.text_input("Enter access passphrase", type="password")
    if not key:
        st.stop()
    if key != secret:
        st.error("Incorrect passphrase.")
        st.stop()

_auth_gate()

# --------------- WRITABLE ROOT -----------------------
def _detect_writable_root() -> Path:
    for p in [Path("./out"), Path(os.path.expanduser("~")) / "out", Path(tempfile.gettempdir()) / "ic-licai-out"]:
        try:
            p.mkdir(parents=True, exist_ok=True)
            t = p / ".touch"
            t.write_text("ok", encoding="utf-8")
            t.unlink()
            return p
        except Exception:
            continue
    return Path(tempfile.gettempdir())

OUT_ROOT = _detect_writable_root()


def _ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


def _safe(name: str) -> str:
    return "".join(c for c in (name or "").strip() if c.isalnum() or c in (" ", "_", "-", ".")).strip().replace(" ", "_")


def _export_bytes(title: str, body: str) -> Tuple[bytes, str, str]:
    base = _safe(title) or "ICLicAI_Report"
    if HAVE_DOCX:
        doc = Document()
        if not PUBLIC_MODE:
            doc.add_paragraph().add_run("CONFIDENTIAL — Internal Evaluation Draft (No Distribution)").bold = True
        doc.add_heading(title, 0)
        for para in body.split("\n\n"):
            doc.add_paragraph(para)
        bio = io.BytesIO()
        doc.save(bio)
        return (
            bio.getvalue(),
            f"{base}.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    if not PUBLIC_MODE:
        body = "CONFIDENTIAL — Internal Evaluation Draft (No Distribution)\n\n" + body
    return body.encode("utf-8"), f"{base}.txt", "text/plain"


def _save_bytes(folder: Path, name: str, data: bytes) -> Tuple[Optional[Path], str]:
    if PUBLIC_MODE:
        return None, "Public mode: server save disabled (download only)."
    try:
        _ensure_dir(folder)
        p = folder / name
        p.write_bytes(data)
        return p, f"Saved to {p}"
    except Exception as e:
        return None, f"Server save skipped ({type(e).__name__}: {e}). Download still works."

# --------------- EVIDENCE EXTRACTION -----------------
TEXT_EXT = {".txt"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
CSV_EXT = {".csv"}
PDF_EXT = {".pdf"}  # filename cue only (kept for future)


def _extract_text_docx(data: bytes) -> str:
    if not HAVE_DOCX:
        return ""
    try:
        bio = io.BytesIO(data)
        doc = Document(bio)
        parts: List[str] = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if txt:
                parts.append(txt)
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                line = " | ".join((cell.text or "").strip() for cell in row.cells)
                if line.strip():
                    parts.append(line)
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_text_pptx(data: bytes) -> str:
    if not HAVE_PPTX:
        return ""
    try:
        bio = io.BytesIO(data)
        prs = Presentation(bio)
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = (shape.text or "").strip()
                    if txt:
                        parts.append(txt)
            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                nt = (slide.notes_slide.notes_text_frame.text or "").strip()
                if nt:
                    parts.append(nt)
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_text_csv(raw: bytes, name: str) -> str:
    """
    CSV semantic extraction: surface headers + a few rows so SME/ESG words
    like 'safety', 'policy', 'contract', 'governance', 'emissions', etc.
    are visible to the heuristics.
    """
    try:
        decoded = raw.decode("utf-8", errors="ignore")
        rows = list(csv.reader(decoded.splitlines()))
        if not rows:
            return ""
        headers = [h.strip() for h in rows[0] if h.strip()]
        cells: List[str] = []
        for row in rows[1:11]:
            cells.extend([c.strip() for c in row if c.strip()])
        header_txt = ", ".join(headers)
        cells_txt = "; ".join(cells)
        return f"CSV:{name}\nHeaders: {header_txt}\nRows: {cells_txt}"
    except Exception:
        try:
            return raw.decode("utf-8", errors="ignore")
        except Exception:
            return ""


def _read_text(files: List[Any]) -> Tuple[str, Dict[str, int], Dict[str, float]]:
    """
    Returns (combined_text, counts_by_ext, weights_used)
    Weights depend on artefact type (contract/JV > SOP/KMP > specs/slides > culture).
    """
    chunks: List[str] = []
    counts: Dict[str, int] = {}
    weights_used: Dict[str, float] = {}

    NAME_WEIGHTS: List[Tuple[str, float]] = [
        ("contract", 1.0),
        ("msa", 1.0),
        ("sow", 0.9),
        ("sla", 0.9),
        ("agreement", 0.9),
        ("joint venture", 1.0),
        ("joint_venture", 1.0),
        ("jv", 1.0),
        ("mou", 1.0),
        ("grant", 0.9),
        ("licence", 0.9),
        ("license", 0.9),
        ("register", 0.9),
        ("knowledge_management", 0.8),
        ("kmp", 0.8),
        ("sop", 0.8),
        ("process", 0.8),
        ("safety", 0.8),
        ("protocol", 0.8),
        ("spec", 0.6),
        ("canvas", 0.6),
        ("bmc", 0.6),
        ("slides", 0.6),
        ("deck", 0.6),
        ("board_pack", 0.8),
        ("board", 0.8),
        ("pricing", 0.7),
        ("tariff", 0.7),
        ("dataset", 0.7),
        ("culture", 0.4),
        ("award", 0.4),
    ]
    EXT_DEFAULTS: Dict[str, float] = {
        ".docx": 0.7,
        ".pptx": 0.6,
        ".txt": 0.5,
        ".csv": 0.6,
        ".pdf": 0.4,
    }

    for f in files or []:
        name = getattr(f, "name", "file")
        lower_name = str(name).lower()
        ext = Path(lower_name).suffix or "none"
        counts[ext] = counts.get(ext, 0) + 1

        weight = EXT_DEFAULTS.get(ext, 0.4)
        for cue, w in NAME_WEIGHTS:
            if cue in lower_name:
                weight = max(weight, w)

        weights_used[lower_name] = weight

        try:
            raw = f.read()
            text = ""
            if ext in TEXT_EXT:
                text = raw.decode("utf-8", errors="ignore")
            elif ext in DOCX_EXT:
                text = _extract_text_docx(raw)
            elif ext in PPTX_EXT:
                text = _extract_text_pptx(raw)
            elif ext in CSV_EXT:
                text = _extract_text_csv(raw, name)
            elif ext in PDF_EXT:
                text = f"[[PDF:{name}]]"
            else:
                text = f"[[FILE:{name}]]"

            if text.strip():
                chunks.append(f"\n# {name}\n{text.strip()}\n")
            else:
                chunks.append(f"\n# {name}\n[[NO-TEXT-EXTRACTED]]\n")
        except Exception:
            chunks.append(f"\n# {name}\n[[READ-ERROR]]\n")

    return "\n".join(chunks).strip(), counts, weights_used

# --------------- SME cues / analysis -----------------
FOUR_LEAF_KEYS: Dict[str, List[str]] = {
    "Human": [
        "team",
        "staff",
        "employee",
        "hire",
        "recruit",
        "training",
        "trained",
        "trainer",
        "onboarding",
        "mentor",
        "apprentice",
        "qualification",
        "certified",
        "cpd",
        "skills matrix",
        "safety training",
        "toolbox talk",
        "rota",
    ],
    "Structural": [
        # IAS 38 explicit / documented assets
        "contract",
        "agreement",
        "joint venture",
        "joint_venture",
        "jv",
        "mou",
        "grant",
        "licence",
        "license",
        "nda",
        "non-disclosure",
        "confidentiality",
        "ip register",
        "asset register",
        "register",
        "policy",
        "sop",
        "procedure",
        "process",
        "workflow",
        "protocol",
        "safety protocol",
        "process map",
        "manual",
        "guide",
        "handbook",
        "template",
        "checklist",
        "board pack",
        "board report",
        "pricing sheet",
        "tariff",
        "datasheet",
        "dataset",
        "qms",
        "iso 9001",
        "iso 27001",
        "knowledge base",
        "architecture",
        "crm",  # explicit system
    ],
    "Customer": [
        "client",
        "customer",
        "account",
        "lead",
        "opportunity",
        "pipeline",
        "quote",
        "proposal",
        "purchase order",
        "po",
        "invoice",
        "renewal",
        "retention",
        "distributor",
        "reseller",
        "channel",
        "customer success",
        "subscription",
    ],
    "Strategic Alliance": [
        "partner",
        "partnership",
        "alliance",
        "strategic",
        "mou",
        "joint venture",
        "framework agreement",
        "collaboration",
        "consortium",
        "university",
        "college",
        "council",
        "ngo",
        "integrator",
        "oem",
        "supplier agreement",
        "grant agreement",
        "licensor",
        "licensee",
        "jv",
    ],
}

TEN_STEPS = [
    "Identify",
    "Separate",
    "Protect",
    "Safeguard",
    "Manage",
    "Control",
    "Use",
    "Monitor",
    "Value",
    "Report",
]

SECTOR_CUES = {
    "GreenTech": [
        "recycling",
        "recycled",
        "waste",
        "biomass",
        "circular",
        "emissions",
        "co2e",
        "solar",
        "pv",
        "turbine",
        "kwh",
        "energy efficiency",
        "retrofit",
        "heat pump",
        "iso 14001",
        "esg",
        "sdg",
        "ofgem",
        "lca",
    ],
    "MedTech": [
        "iso 13485",
        "mhra",
        "ce mark",
        "clinical",
        "gcp",
        "patient",
        "medical device",
        "pms",
        "post-market surveillance",
    ],
    "AgriTech": [
        "soil",
        "irrigation",
        "seed",
        "fertiliser",
        "biomass",
        "yield",
        "farm",
    ],
}

# Explicit structural cues (IAS 38-compliant artefact hints)
EXPLICIT_STRUCTURAL_CUES: List[str] = [
    "contract",
    "agreement",
    "msa",
    "sow",
    "sla",
    "mou",
    "joint venture",
    "joint_venture",
    "jv",
    "grant",
    "licence",
    "license",
    "ip register",
    "asset register",
    "register",
    "policy",
    "sop",
    "procedure",
    "process",
    "workflow",
    "protocol",
    "safety protocol",
    "process map",
    "manual",
    "guide",
    "handbook",
    "template",
    "checklist",
    "board pack",
    "board report",
    "pricing sheet",
    "tariff",
    "datasheet",
    "dataset",
    "qms",
    "iso 9001",
    "iso 27001",
    "crm",
]

# ESG & Seven Stakeholder cues (Sugai / Weir)
ESG_CUES: List[str] = [
    "esg",
    "sdg",
    "carbon",
    "emissions",
    "net zero",
    "scope 1",
    "scope 2",
    "scope 3",
    "sustainability",
    "governance",
    "diversity",
    "inclusion",
    "impact",
    "stakeholder",
]

SEVEN_STAKEHOLDER_CUES: List[str] = [
    "employee",
    "staff",
    "worker",
    "investor",
    "shareholder",
    "lender",
    "customer",
    "client",
    "supplier",
    "vendor",
    "partner",
    "alliance",
    "community",
    "local authority",
    "municipality",
    "ngo",
    "nature",
    "environment",
    "biodiversity",
]

# --------------- ANALYSIS ENGINE ---------------------
def _analyse_weighted(
    text: str,
    weights_by_file: Dict[str, float],
) -> Tuple[Dict[str, Any], Dict[str, float], Dict[str, Any], int]:
    """
    Weighted Four-Leaf & Ten-Steps.
    Returns:
      ic_map (with tick/narrative/score),
      leaf_scores (raw weighted scores for 4-leaf),
      ten (scores+narratives),
      quality% (heuristic)
    """
    sector = st.session_state.get("sector", "Other")
    t_all = (text or "").lower()

    leaf_scores: Dict[str, float] = {
        "Human": 0.0,
        "Structural": 0.0,
        "Customer": 0.0,
        "Strategic Alliance": 0.0,
    }
    step_scores: Dict[str, float] = {s: 0.0 for s in TEN_STEPS}

    sector_present = False
    if sector in SECTOR_CUES:
        if any(c in t_all for c in SECTOR_CUES[sector]):
            sector_present = True

    # ----- Structural vs Tacit weighting -----
    max_weight = max(weights_by_file.values() or [0.4])

    # Base structural emphasis from explicit cues anywhere in the text (IAS 38 explicit assets)
    for cue in EXPLICIT_STRUCTURAL_CUES:
        if cue in t_all:
            leaf_scores["Structural"] += max_weight * 1.5  # audit-ready bump

    # Four-Leaf cues (with sector reinforcement)
    for leaf, cues in FOUR_LEAF_KEYS.items():
        eff = list(cues)
        if sector_present and leaf in ("Structural", "Customer", "Strategic Alliance"):
            eff += SECTOR_CUES[sector]
        base = 0.0
        for cue in eff:
            if cue in t_all:
                base += max_weight
        leaf_scores[leaf] += base

    # Ten-Steps scoring (file-name based + ESG / FRAND cues)
    def bump(step: str, amt: float) -> None:
        step_scores[step] = step_scores.get(step, 0.0) + amt

    for fname, w in (weights_by_file or {}).items():
        n = fname.lower()

        # Contracts / grants / agreements → Structural dominates, Customer/SA secondary
        if any(k in n for k in ["contract", "msa", "sow", "sla", "po", "agreement"]):
            leaf_scores["Structural"] += 2.5 * w
            leaf_scores["Customer"] += 1.0 * w
            bump("Control", 2.0 * w)
            bump("Use", 2.5 * w)

        if any(k in n for k in ["joint_venture", "joint venture", "jv", "mou", "grant"]):
            leaf_scores["Structural"] += 2.5 * w
            leaf_scores["Strategic Alliance"] += 1.5 * w
            bump("Control", 2.0 * w)
            bump("Use", 2.0 * w)

        # Knowledge / SOP / KMP / safety / ISO → Structural + Human
        if any(k in n for k in ["knowledge", "kmp", "sop", "process", "safety", "protocol", "risk", "qms", "iso"]):
            leaf_scores["Structural"] += 1.8 * w
            leaf_scores["Human"] += 0.8 * w
            bump("Identify", 1.8 * w)
            bump("Separate", 1.4 * w)
            bump("Manage", 1.6 * w)
            bump("Safeguard", 1.0 * w)

        # Specs/slides/canvas → Structural + Use
        if any(k in n for k in ["spec", "canvas", "deck", "slides", "pptx"]):
            leaf_scores["Structural"] += 0.8 * w
            bump("Identify", 0.8 * w)
            bump("Use", 0.6 * w)

        # Pricing/licensing hints → Use/Value (multi value streams)
        if any(k in n for k in ["price", "pricing", "royalty", "subscription", "oem", "white label"]):
            bump("Use", 1.2 * w)
            bump("Value", 1.6 * w)

        # Governance/reporting → Monitor/Report
        if any(k in n for k in ["board", "report", "dashboard", "audit"]):
            bump("Report", 1.4 * w)
            bump("Monitor", 1.2 * w)

    # ESG & Seven Stakeholder presence → boost Report/Value (double materiality)
    esg_hits = any(c in t_all for c in ESG_CUES)
    stakeholder_hits = any(c in t_all for c in SEVEN_STAKEHOLDER_CUES)
    if esg_hits or stakeholder_hits:
        bump("Report", 1.2)
        bump("Value", 1.0)

    if sector_present:
        bump("Use", 0.8)
        bump("Report", 0.5)

    # Make sure Structural "wins" when explicit + tacit both present:
    # if Structural>0 and (Customer or SA also high), add a small dominance bump.
    if leaf_scores["Structural"] > 0 and (leaf_scores["Customer"] > 0 or leaf_scores["Strategic Alliance"] > 0):
        leaf_scores["Structural"] *= 1.15  # dominance tweak

    # Convert leaf_scores -> ticks & narratives
    ic_map: Dict[str, Any] = {}
    avg_leaf = (sum(leaf_scores.values()) / max(1, len(leaf_scores)))
    threshold = max(1.0, avg_leaf * 0.6)

    for leaf, score in leaf_scores.items():
        tick = score >= threshold
        if leaf == "Human":
            nar = (
                "Human Capital evidenced (values, awards, training and safety practice), but competency and role-mapping "
                "should be consolidated into a formal skills register."
                if tick
                else "Human Capital is not yet clearly evidenced; competency mapping, training logs and safety records are needed."
            )
        elif leaf == "Structural":
            nar = (
                "Structural Capital appears IAS 38-ready in places (contracts, SOPs, protocols, registers, CRM and board packs "
                "are present), supporting audit-ready recognition on the balance sheet."
                if tick
                else "Structural Capital is under-documented; explicit artefacts (contracts, registers, SOPs, board packs, "
                "pricing, datasets, CRM) should be consolidated into an auditable IA Register."
            )
        elif leaf == "Customer":
            nar = (
                "Customer Capital is evidenced through relationships, renewal logic and channels, supporting recurring value capture "
                "and future licensing opportunities."
                if tick
                else "Customer Capital appears weak in the evidence; relationship histories, renewals, CRM and pipeline data should be structured."
            )
        else:
            nar = (
                "Strategic Alliance Capital is evidenced (JVs, MoUs, partners, universities, councils), enabling co-creation "
                "and ecosystem-based licensing opportunities."
                if tick
                else "Strategic alliances are not clearly evidenced; JV/MoU documentation and partner frameworks are needed."
            )
        ic_map[leaf] = {"tick": tick, "narrative": nar, "score": round(score, 2)}

    # Ten-Steps scores
    base = 3.0
    ten_scores: List[int] = []
    ten_narrs: List[str] = []

    for step in TEN_STEPS:
        s_float = base + step_scores.get(step, 0.0)
        s = int(max(1, min(10, round(s_float))))
        ten_scores.append(s)
        ten_narrs.append(f"{step}: readiness ≈ {s}/10.")

    ten = {"scores": ten_scores, "narratives": ten_narrs}

    # Evidence quality metric
    files_factor = min(1.0, len(weights_by_file) / 6.0)
    leaf_div = sum(1 for v in ic_map.values() if v["tick"]) / 4.0
    weight_mean = (sum(weights_by_file.values()) / max(1, len(weights_by_file))) if weights_by_file else 0.4
    quality = int(round(100 * (0.45 * files_factor + 0.35 * leaf_div + 0.20 * min(1.0, weight_mean))))

    return ic_map, leaf_scores, ten, quality

# --------------- INTERPRETIVE NARRATIVE --------------
def _build_interpreted_summary(
    case: str,
    leaf_scores: Dict[str, float],
    ic_map: Dict[str, Any],
    ten: Dict[str, Any],
    evidence_quality: int,
    context: Dict[str, str],
) -> str:
    sector = st.session_state.get("sector", "Other")
    size = st.session_state.get("company_size", "Micro (1–10)")

    strengths = [k for k, v in ic_map.items() if v.get("tick")]
    gaps = [k for k, v in ic_map.items() if not v.get("tick")]

    ts = ten.get("scores") or [5] * len(TEN_STEPS)
    strong_steps = [s for s, sc in zip(TEN_STEPS, ts) if sc >= 7]
    weak_steps = [s for s, sc in zip(TEN_STEPS, ts) if sc <= 5]

    # Detect whether ESG & Seven Stakeholder cues are present
    narrative_text = context.get("why", "") + " " + context.get("markets", "")
    seven_hit = any(c in narrative_text.lower() for c in SEVEN_STAKEHOLDER_CUES)
    esg_hit = any(c in narrative_text.lower() for c in ESG_CUES)

    # 1) Context & positioning
    p1 = (
        f"{case} is a {size} in {sector}. Based on uploaded artefacts and company context, the company shows an "
        f"emerging ability to codify and scale its operating model, with measurable signals across "
        f"{', '.join(strengths) if strengths else 'selected IC dimensions'}."
    )

    # 2) Four-Leaf interpretation (explicit vs tacit, IAS 38)
    if strengths:
        p2a = f"Strengths concentrate in {', '.join(strengths)}" + (f"; gaps are {', '.join(gaps)}." if gaps else ".")
    else:
        p2a = "Strengths are not yet well-evidenced; additional artefacts are required."

    p2b = (
        "Evidence points to maturing Structural Capital where explicit artefacts — contracts, SOPs, protocols, registers, "
        "board materials, CRM and datasets — are present. These are the primary candidates for IAS 38-compliant recognition on "
        "the balance sheet. Human, Customer and Strategic Alliance Capital are reflected mainly through tacit know-how, "
        "relationships and informal practice, which require codification before they become audit-ready."
    )
    p2 = p2a + " " + p2b

    # 3) Ten-Steps insight and readiness for licensing
    if strong_steps or weak_steps:
        p3 = (
            f"Ten-Steps patterns indicate strong {', '.join(strong_steps) if strong_steps else 'foundations'}; "
            f"progress is constrained by {', '.join(weak_steps) if weak_steps else 'later-stage governance, valuation and reporting readiness'}."
        )
    else:
        p3 = (
            "Ten-Steps scores suggest a developing baseline; company-side review will refine scoring as artefacts are "
            "consolidated and IA governance is embedded."
        )

    # 4) Seven Stakeholder / ESG framing and value streams
    if seven_hit or esg_hit:
        p4_intro = (
            "Using the Seven Stakeholder Model (SSM) as defined by Professor Philip Sugai and Dr Maria Weir, "
            "the company can frame value creation across employees, investors, customers, partners and suppliers, "
            "communities and the natural environment. This provides a structured way to connect ESG performance "
            "and double materiality to concrete intangible assets."
        )
    else:
        p4_intro = (
            "Applying the Seven Stakeholder Model (SSM) as defined by Professor Philip Sugai and Dr Maria Weir "
            "would allow the company to frame value creation across employees, investors, customers, partners and "
            "suppliers, communities and the natural environment, even if these links are not yet fully articulated "
            "in the evidence set."
        )

    p4_mid = (
        "From a commercialisation perspective, explicit Structural Capital (contracts, data, software, methods, indices, "
        "protocols) can support multiple simultaneous value streams — including revenue licences, access or community "
        "licences, co-creation arrangements and data/algorithm sharing — provided that ownership, rights and governance "
        "are clarified."
    )

    actions = [
        "Create a single IA Register linking all explicit artefacts (contracts, JVs, SOPs, protocols, datasets, board packs, CRM) to the 4-Leaf Model.",
        "Map each explicit asset to at least one licensing-ready value stream (revenue, access/community, co-creation, defensive or data/algorithm sharing).",
        "Introduce quarterly governance reporting (board pack + KPI dashboard) to strengthen Monitor and Report and to evidence ESG and stakeholder impacts.",
        "Define valuation approach (IAS 38 fair value) and link to licensing templates so that audit-ready Structural Capital supports near-term monetisation.",
        "Formalise competency matrices and training logs so that tacit Human Capital can be progressively codified into Structural Capital.",
    ]
    p4_actions = "Assumptions & Action Plan:\n" + "\n".join([f"• {a}" for a in actions])

    p4 = p4_intro + "\n\n" + p4_mid + "\n\n" + p4_actions

    # 5) Evidence quality and next evidence requests
    missing = (
        "Request additional artefacts: CRM/renewal data, NDA/licence/royalty and access terms, IA or IP registers, "
        "board/management reports, and any ESG or stakeholder dashboards used in internal decision-making."
    )
    p5 = f"Evidence quality ≈ {evidence_quality}% (heuristic). {missing}"

    return "\n\n".join([p1, p2, p3, p4, p5])

# --------- COMPANY CONTEXT AUTO-SPLIT HELPER ----------
def _auto_split_expert_block(text: str) -> Dict[str, str]:
    """
    Take a single pasted block and try to split it across:
    why_service, stage, plan_s, plan_m, plan_l, markets_why, sale_price_why
    using blank lines or sentence boundaries.
    """
    t = (text or "").strip()
    keys = [
        "why_service",
        "stage",
        "plan_s",
        "plan_m",
        "plan_l",
        "markets_why",
        "sale_price_why",
    ]
    if not t:
        return {k: "" for k in keys}

    blocks = [b.strip() for b in t.replace("\r\n", "\n").split("\n\n") if b.strip()]

    if len(blocks) < 5:
        blocks = [s.strip() for s in re.split(r"(?<=[.!?])\s+", t) if s.strip()]

    out: Dict[str, str] = {k: "" for k in keys}
    for k, chunk in zip(keys, blocks):
        out[k] = chunk
    return out

# ------------------ SESSION DEFAULTS -----------------
ss = st.session_state
ss.setdefault("case_name", "Untitled Company")
ss.setdefault("company_size", "Micro (1–10)")
ss.setdefault("company_type", "Registered company / SME")
ss.setdefault("sector", "Other")
ss.setdefault("uploads", [])
ss.setdefault("combined_text", "")
ss.setdefault("ic_map", {})
ss.setdefault("ten_steps", {})
ss.setdefault("narrative", "")
ss.setdefault("leaf_scores", {})
ss.setdefault("evidence_quality", 0)
ss.setdefault("file_counts", {})

# Expert prompts
ss.setdefault("why_service", "")
ss.setdefault("stage", "")
ss.setdefault("plan_s", "")
ss.setdefault("plan_m", "")
ss.setdefault("plan_l", "")
ss.setdefault("markets_why", "")
ss.setdefault("sale_price_why", "")
ss.setdefault("full_context_block", "")
ss.setdefault("auto_split_on_save", True)

# LIP Assistant state
ss.setdefault("lip_history", [])

# Asset verification (overall notes)
ss.setdefault("verification_notes", "")

SIZES = [
    "Micro (1–10)",
    "Small (11–50)",
    "Medium (51–250)",
    "Large (250+)",
]

SECTORS = [
    "Food & Beverage",
    "MedTech",
    "GreenTech",
    "AgriTech",
    "Biotech",
    "Software/SaaS",
    "FinTech",
    "EdTech",
    "Manufacturing",
    "Creative/Digital",
    "Professional Services",
    "Mobility/Transport",
    "Energy",
    "Other",
]

# New: organisation / project type (covers pre-company & spin-outs)
COMPANY_TYPES = [
    "Registered company / SME",
    "Large corporate / group",
    "Pre-startup (idea or prototype)",
    "Spin-out from an existing organisation",
    "Innovation hub / lab",
    "Research-based (university or institute)",
    "Other",
]

## --------------- SIDEBAR BRANDING & NAV ---------------
with st.sidebar:
    # IMPAC3T-IP logo (top) – safe load, moderate width
    try:
        if Path(IMPACT3T_LOGO_PATH).is_file():
            st.image(IMPACT3T_LOGO_PATH, width=170)
        else:
            st.markdown("**IMPAC3T-IP**")
    except Exception:
        st.markdown("**IMPAC3T-IP**")

    st.markdown("---")

    # Navigation (middle of sidebar)
    st.markdown("### Navigate")
    page = st.radio(
        "Go to:",
        [
            "Company",
            "Analyse Evidence",
            "Asset Verification",
            "LIP Console",
            "Reports",
            "Licensing Templates",
            "LIP Assistant",
            "Glossary",
        ],
        index=0,
    )

    st.markdown("---")

    # EU flag + funding line – **footer**
    try:
        if Path(EU_FLAG_PATH).is_file():
            st.image(EU_FLAG_PATH, width=80)
        else:
            st.markdown("EU-funded tool")
    except Exception:
        st.markdown("EU-funded tool")

    st.caption(
        "This tool has been developed within the IMPAC3T-IP project, which has received "
        "funding from the European Union's Horizon Europe programme under Grant Agreement "
        "No. 101135832."
    )
# -------------------- PAGES -------------------------

# 1) COMPANY (licensing-focused SME questions with dropdowns)
if page == "Company":
    st.header("Company / project details")

    with st.form("company_form"):
        c1, c2, c3 = st.columns([1.1, 1, 1])
        with c1:
            case_name = st.text_input(
                "Company or project name *",
                ss.get("case_name", ""),
                help="If you don’t have a registered company yet, use your project or working title."
            )
        with c2:
            size = st.selectbox(
                "Size (now or planned)",
                SIZES,
                index=SIZES.index(ss.get("company_size", SIZES[0])),
            )
        with c3:
            current_sector = ss.get("sector", "Other")
            sector_index = SECTORS.index(current_sector) if current_sector in SECTORS else SECTORS.index("Other")
            sector = st.selectbox(
                "Sector / Industry",
                SECTORS,
                index=sector_index,
            )

        # NEW: organisation / project type (covers pre-company & spin-outs)
        stored_type = ss.get("company_type", "Registered company / SME")
        if stored_type in COMPANY_TYPES:
            company_type_index = COMPANY_TYPES.index(stored_type)
        else:
            company_type_index = 0

        company_type = st.selectbox(
            "What best describes this organisation or project?",
            COMPANY_TYPES,
            index=company_type_index,
            help="Covers registered companies, pre-startups, spin-outs, innovation hubs, university-based projects and large corporates.",
        ~)

        st.markdown("#### Simple questions to set the scene (for licensing)")

      st.markdown("#### Valuation & Exit Timing (for planning)")

colv1, colv2 = st.columns(2)

import datetime

st.markdown("#### Valuation & Exit Timing (for planning)")

colv1, colv2 = st.columns(2)

with colv1:
    last_valuation_date = st.date_input(
        "Date of last valuation (if any)",
        value=datetime.date(2025, 12, 31),
        key="last_valuation_date",
        help="If a valuation has already been carried out, enter the date. Leave blank if not applicable."
    )

with colv2:
    exit_date = st.date_input(
        "Suggested or expected exit date",
        value=datetime.date(2027, 1, 1),
        key="exit_date",
        help="When do you expect a sale, investment, or restructuring event to take place?"
    )
    )

        # --- Q1: What are you working on? ---
        q1_type_options = [
            "New product",
            "New service",
            "Software / app",
            "Online platform or portal",
            "Training or learning content",
            "Data or analytics",
            "Method / process or toolkit",
            "Brand / marketing concept",
            "Other idea",
        ]
        q1_type = st.multiselect(
            "1) What are you working on? *",
            q1_type_options,
            help="Tick all that apply. This can be a company, project, spin-out, or early idea.",
        )
        q1_desc = st.text_area(
            "Short description (one or two sentences)",
            ss.get("why_service", ""),
            height=70,
        )

        # --- Q2: Where are you on the journey? ---
        q2_stage_options = [
            "Idea only / early concept",
            "Prototype or proof-of-concept",
            "Pilot with first users",
            "First paying customers",
            "Growing / scaling",
            "Established product or service",
        ]
        q2_stage_choice = st.selectbox(
            "2) Where are you on the journey right now? *",
            q2_stage_options,
            index=0,
            help="Roughly where you are today — this doesn’t need to be perfect.",
        )
        q2_notes = st.text_area(
            "Anything important about your current stage (optional)",
            ss.get("stage", ""),
            height=60,
        )

        # --- Q3: Who should use this, and where? ---
        q3_user_options = [
            "Households / general public",
            "Small businesses / SMEs",
            "Large companies",
            "Government / public bodies",
            "Hospitals / clinics",
            "Schools / universities",
            "Farmers or producers",
            "NGOs / community groups",
            "Other (not listed)",
        ]
        q3_region_options = [
            "Local only",
            "National",
            "Regional (e.g. EU, East Africa)",
            "Across Africa",
            "Across Europe",
            "Global",
        ]
        q3_users = st.multiselect(
            "3) Who do you want to use this? *",
            q3_user_options,
            help="Tick all that apply.",
        )
        q3_regions = st.multiselect(
            "Where do you mainly want to use or sell it? *",
            q3_region_options,
        )
        q3_notes = st.text_area(
            "Key countries or regions (optional)",
            ss.get("plan_s", ""),
            height=60,
        )

        # --- Q4: What do you already have written down or built? ---
        q4_asset_options = [
            "Nothing written down yet",
            "Notes or concept document",
            "Slides / pitch deck",
            "Business plan or canvas",
            "Prototype / demo",
            "Working software / app",
            "Datasets or analytics",
            "Training or learning materials",
            "Brand assets (name, logo, style)",
            "Website or landing page",
            "Policies, SOPs or manuals",
            "Contracts or agreements",
        ]
        q4_assets = st.multiselect(
            "4) What do you already have written down or built? *",
            q4_asset_options,
            help="Tick everything you already have. The documents themselves can be uploaded below.",
        )
        q4_notes = st.text_area(
            "Anything else you have already created (optional)",
            ss.get("plan_m", ""),
            height=60,
        )

        # --- Q5: Who else is involved, and what’s agreed? ---
        q5_who_options = [
            "Just me / us (founder team)",
            "Co-founders",
            "University or research institute",
            "Current employer",
            "Investor or funder",
            "Customer / client",
            "Supplier or tech partner",
            "Government or public body",
            "NGO / community group",
            "Other partner",
        ]
        q5_agree_options = [
            "Nothing formal yet",
            "Informal discussions only",
            "Emails that talk about roles or rights",
            "NDA / confidentiality agreement",
            "Grant agreement",
            "Commercial contract",
            "IP or licence agreement",
            "Don’t know / need to check",
        ]
        q5_who = st.multiselect(
            "5) Who else is involved? *",
            q5_who_options,
        )
        q5_agreed = st.multiselect(
            "What (if anything) has already been agreed in writing? *",
            q5_agree_options,
        )
        q5_notes = st.text_area(
            "Anything sensitive or important about these relationships (optional)",
            ss.get("plan_l", ""),
            height=60,
        )

        # --- Q6: How do you hope to earn from this, and what are you happy to share? ---
        q6_earn_options = [
            "One-off sales",
            "Subscription (monthly or yearly)",
            "Pay-per-use",
            "Revenue share",
            "Licence fees / royalties",
            "Consulting or services",
            "Advertising / sponsorship",
            "Not sure yet",
        ]
        q6_share_options = [
            "Keep fully private / closed",
            "Only for paying customers",
            "Lower-cost access for SMEs / start-ups",
            "Free or low-cost for community / public sector",
            "Open for non-commercial use",
            "Open source (code or content)",
            "Not sure yet",
        ]
        q6_earn_sel = st.multiselect(
            "6) How do you hope to earn from this? *",
            q6_earn_options,
        )
        q6_share_sel = st.multiselect(
            "What are you happy to share or make easier to access? *",
            q6_share_options,
        )
        q6_notes = st.text_area(
            "Anything else about pricing, access or fairness (optional)",
            ss.get("markets_why", ""),
            height=60,
        )

        st.caption("Uploads are held in session until analysis. Nothing is written to server until export.")
        uploads = st.file_uploader(
            "Upload evidence (PDF, DOCX, TXT, CSV, XLSX, PPTX, images)",
            type=["pdf", "docx", "txt", "csv", "xlsx", "pptx", "png", "jpg", "jpeg", "webp"],
            accept_multiple_files=True,
            key="uploader_main",
        )

        submitted = st.form_submit_button("Save details")

        if submitted:
            # Build simple text answers from dropdowns + notes
            # Q1
            q1_parts = []
            if q1_type:
                q1_parts.append("Types: " + ", ".join(q1_type))
            if q1_desc.strip():
                q1_parts.append("Description: " + q1_desc.strip())
            q1_what = " ".join(q1_parts).strip()

            # Q2
            q2_parts = [q2_stage_choice]
            if q2_notes.strip():
                q2_parts.append(q2_notes.strip())
            q2_text = " ".join(q2_parts).strip()

            # Q3
            q3_parts = []
            if q3_users:
                q3_parts.append("Users: " + ", ".join(q3_users))
            if q3_regions:
                q3_parts.append("Regions: " + ", ".join(q3_regions))
            if q3_notes.strip():
                q3_parts.append("Notes: " + q3_notes.strip())
            q3_text = " ".join(q3_parts).strip()

            # Q4
            q4_parts = []
            if q4_assets:
                q4_parts.append("Assets already in place: " + ", ".join(q4_assets))
            if q4_notes.strip():
                q4_parts.append("Other: " + q4_notes.strip())
            q4_text = " ".join(q4_parts).strip()

            # Q5
            q5_parts = []
            if q5_who:
                q5_parts.append("People / organisations involved: " + ", ".join(q5_who))
            if q5_agreed:
                q5_parts.append("What’s agreed: " + ", ".join(q5_agreed))
            if q5_notes.strip():
                q5_parts.append("Notes: " + q5_notes.strip())
            q5_text = " ".join(q5_parts).strip()

            # Q6
            q6_parts = []
            if q6_earn_sel:
                q6_parts.append("Revenue ideas: " + ", ".join(q6_earn_sel))
            if q6_share_sel:
                q6_parts.append("Sharing / access: " + ", ".join(q6_share_sel))
            if q6_notes.strip():
                q6_parts.append("Notes: " + q6_notes.strip())
            q6_text = " ".join(q6_parts).strip()

            # Required fields check (six questions + name)
            missing = [
                ("Company or project name", case_name),
                ("What are you working on?", q1_what),
                ("Where are you on the journey?", q2_text),
                ("Who should use this, and where?", q3_text),
                ("What you already have written down or built", q4_text),
                ("Who else is involved, and what’s agreed", q5_text),
                ("How you hope to earn from this, and what you’re happy to share", q6_text),
            ]
            missing_fields = [label for (label, val) in missing if not (val or "").strip()]

            if missing_fields:
                st.error("Please complete required fields: " + ", ".join(missing_fields))
            else:
                # Store everything back into session (reuse existing keys)
                ss["case_name"] = case_name
                ss["company_size"] = size
                ss["company_type"] = company_type
                ss["sector"] = sector

                # Map to existing keys used by the analysis engine
                ss["why_service"] = q1_what
                ss["stage"] = q2_text
                ss["plan_s"] = q3_text
                ss["plan_m"] = q4_text
                ss["plan_l"] = q5_text
                ss["markets_why"] = q6_text
                # sale_price_why stays as-is / unused in this simplified licensing view

                # Clear legacy “full story” field (no longer used)
                ss["full_context_block"] = ""
                ss["auto_split_on_save"] = False

                if uploads:
                    ss["uploads"] = uploads

                st.success("Saved company / project details and licensing context.")

    if ss.get("uploads"):
        st.info(f"{len(ss['uploads'])} file(s) stored in session. Go to **Analyse Evidence** next.")
        
# 2) ANALYSE EVIDENCE (with radar / evidence quality)
elif page == "Analyse Evidence":
    st.header("Evidence Dashboard & Analysis")

    col1, col2 = st.columns([1, 1])
    with col1:
        st.subheader("Evidence Quality")
        eq = int(ss.get("evidence_quality", 0))
        st.progress(min(100, max(0, eq)) / 100.0)
        st.caption(f"{eq}% coverage (heuristic — based on artefact mix and IC diversity).")

        counts = ss.get("file_counts", {}) or {}
        if counts:
            st.markdown("**Files by type (session):**")
            for ext, n in counts.items():
                st.markdown(f"- `{ext}` → {n} file(s)")
        else:
            st.caption("No files analysed yet.")

    with col2:
        st.subheader("IC Radar (4-Leaf + Ten-Steps)")

        ic_map: Dict[str, Any] = ss.get("ic_map", {})
        ten = ss.get(
            "ten_steps",
            {"scores": [5] * len(TEN_STEPS), "narratives": [f"{s}: tbd" for s in TEN_STEPS]},
        )

        leaf_labels = ["Human", "Structural", "Customer", "Strategic Alliance"]
        leaf_vals = [float(ic_map.get(l, {}).get("score", 0.0)) for l in leaf_labels]
        if any(v > 0 for v in leaf_vals):
            fig_leaf = go.Figure()
            fig_leaf.add_trace(
                go.Scatterpolar(
                    r=leaf_vals + leaf_vals[:1],
                    theta=leaf_labels + leaf_labels[:1],
                    fill="toself",
                    name="IC Intensity",
                )
            )
            fig_leaf.update_layout(
                polar=dict(radialaxis=dict(visible=True, range=[0, max(leaf_vals) or 1])),
                showlegend=False,
                margin=dict(l=20, r=20, t=20, b=20),
            )
            st.plotly_chart(fig_leaf, use_container_width=True)
        else:
            st.caption("Radar will appear once analysis has been run and IC signals are detected.")

        step_scores = ten.get("scores") or [5] * len(TEN_STEPS)
        if step_scores:
            fig_steps = go.Figure(
                data=[
                    go.Bar(
                        x=TEN_STEPS,
                        y=step_scores,
                    )
                ]
            )
            fig_steps.update_layout(
                yaxis=dict(range=[0, 10]),
                margin=dict(l=20, r=20, t=20, b=40),
            )
            st.plotly_chart(fig_steps, use_container_width=True)

    st.markdown("---")
    st.subheader("Interpreted Analysis Summary (first 5000 chars)")
    combined = ss.get("combined_text", "")
    st.text_area(
        "Summary view (read-only; editable in LIP Console)",
        combined[:5000],
        height=260,
        key="combined_preview",
    )

    if st.button("Run analysis now"):
        uploads: List[Any] = ss.get("uploads") or []
        extracted, counts, weights = _read_text(uploads)

        ss["file_counts"] = counts or {}

        context = {
            "why": ss.get("why_service", ""),
            "stage": ss.get("stage", ""),
            "plan_s": ss.get("plan_s", ""),
            "plan_m": ss.get("plan_m", ""),
            "plan_l": ss.get("plan_l", ""),
            "markets": ss.get("markets_why", ""),
            "sale": ss.get("sale_price_why", ""),
        }

        context_stub = (
            f"[CTX] why={context['why'][:140]} | stage={context['stage'][:140]} | "
            f"plans=({context['plan_s'][:60]}/{context['plan_m'][:60]}/{context['plan_l'][:60]}) | "
            f"markets={context['markets'][:140]} | sale={context['sale'][:60]}"
        )

        combined_text_for_detection = (extracted + "\n\n" + context_stub).strip().lower()

        ic_map, leaf_scores, ten, quality = _analyse_weighted(combined_text_for_detection, weights)

        case = ss.get("case_name", "Untitled Company")
        interpreted = _build_interpreted_summary(case, leaf_scores, ic_map, ten, quality, context)

        ss["combined_text"] = interpreted
        ss["ic_map"] = ic_map
        ss["ten_steps"] = ten
        ss["leaf_scores"] = leaf_scores
        ss["evidence_quality"] = quality

        if len(extracted.strip()) < 100:
            st.warning(
                "Little machine-readable text was extracted (DOCX/PPTX/CSV extraction is enabled). "
                "If PDFs dominate, consider adding a brief TXT note or export key pages to DOCX."
            )

        st.success("Analysis complete. Open **LIP Console** to refine and export.")
# 3) ASSET VERIFICATION (human check of assets & ESG claims)
elif page == "Asset Verification":
    st.header("Asset Verification — Evidence & ESG checks")

    uploads = ss.get("uploads") or []

    if not uploads:
        st.warning("No evidence files found in session. Go to **Company** and upload documents first.")
    else:
        st.markdown(
            "This page supports a **human verification step** for the documents already uploaded. "
            "It helps the Licensing & Intangibles Partner (LIP) or TTO to check ownership, validity "
            "and the robustness of ESG or impact claims before moving to the LIP Console and reports."
        )

        st.info(
            "Use this view to **challenge the evidence**: confirm who owns what, whether contracts are in force, "
            "and whether ESG or impact language is backed by real proof (not greenwashing)."
        )

        OPTION_LABELS = ["Select…", "Yes", "No", "Unsure / needs follow-up"]

        for f in uploads:
            fname = getattr(f, "name", "(unnamed file)")
            safe_key = _safe(fname)

            st.markdown(f"### {fname}")

            col1, col2, col3 = st.columns(3)
            with col1:
                st.selectbox(
                    "Ownership / rights clear?",
                    OPTION_LABELS,
                    key=f"ver_own_{safe_key}",
                    help="Does this document clearly show who owns or controls the asset?",
                )
            with col2:
                st.selectbox(
                    "Up to date & in force?",
                    OPTION_LABELS,
                    key=f"ver_date_{safe_key}",
                    help="Is the document current (still valid, signed, dates make sense)?",
                )
            with col3:
                st.selectbox(
                    "Claims supported by evidence?",
                    OPTION_LABELS,
                    key=f"ver_claim_{safe_key}",
                    help="Where the document makes bold or ESG-related claims, are they backed by real detail?",
                )

            st.text_area(
                "Notes / follow-up for this document",
                key=f"ver_notes_{safe_key}",
                height=80,
                help="Capture anything that needs clarification, extra evidence, or legal review.",
            )

            st.markdown("---")

        st.subheader("Overall verification summary for this company")
        ss["verification_notes"] = st.text_area(
            "Overall verification notes",
            ss.get("verification_notes", ""),
            height=120,
            help="High-level view: which assets look robust, which are weak, and what needs to happen next.",
        )

        st.caption(
            "These checks do **not** replace legal review. They are a structured way for the LIP / TTO to "
            "document how far the evidence can be trusted before licensing design and valuation."
        )
# 4) LIP CONSOLE (was Expert View)
elif page == "LIP Console":
    st.header("LIP Console — Narrative & IC Map")
    nar = st.text_area(
        "Summary (editable by Licensing & Intangibles Partner)",
        value=ss.get("combined_text", ""),
        height=220,
        key="nar_edit",
    )
    ss["narrative"] = nar or ss.get("narrative", "")
    ss["combined_text"] = ss["narrative"]

    colA, colB = st.columns([1, 1])
    with colA:
        if not PUBLIC_MODE:
            st.subheader("Evidence Quality")
            st.progress(min(100, max(0, ss.get("evidence_quality", 0))) / 100.0)
            st.caption(f"{ss.get('evidence_quality', 0)}% evidence coverage (heuristic)")

        st.subheader("4-Leaf Map")
        ic_map: Dict[str, Any] = ss.get("ic_map", {})
        for leaf in ["Human", "Structural", "Customer", "Strategic Alliance"]:
            row = ic_map.get(
                leaf,
                {"tick": False, "narrative": f"No assessment yet for {leaf}.", "score": 0.0},
            )
            tick = "✓" if row["tick"] else "•"
            suffix = "" if PUBLIC_MODE else f"  _(score: {row.get('score', 0.0)})_"
            st.markdown(f"- **{leaf}**: {tick}{suffix}")
            st.caption(row["narrative"])

        st.subheader("Company context (read-only)")
        st.markdown(f"- **Why service:** {ss.get('why_service', '') or '—'}")
        st.markdown(f"- **Stage:** {ss.get('stage', '') or '—'}")
        st.markdown(
            f"- **Plans:** S={ss.get('plan_s', '') or '—'} | "
            f"M={ss.get('plan_m', '') or '—'} | L={ss.get('plan_l', '') or '—'}"
        )
        st.markdown(f"- **Markets & why:** {ss.get('markets_why', '') or '—'}")
        st.markdown(f"- **Target sale & why:** {ss.get('sale_price_why', '') or '—'}")

    with colB:
        st.subheader("Ten-Steps Readiness")
        raw_ten = ss.get("ten_steps") or {}
        scores = raw_ten.get("scores") or [5] * len(TEN_STEPS)
        narrs = raw_ten.get("narratives") or [f"{s}: tbd" for s in TEN_STEPS]
        ten = {"scores": scores, "narratives": narrs}

        st.dataframe(
            {"Step": TEN_STEPS, "Score (1–10)": ten["scores"]},
            hide_index=True,
            use_container_width=True,
        )
        with st.expander("Narrative per step"):
            for s, n in zip(TEN_STEPS, ten["narratives"]):
                st.markdown(f"**{s}** — {n}")

# 5) REPORTS
elif page == "Reports":
    st.header("Reports & Exports")
    case_name = ss.get("case_name", "Untitled Company")
    case_folder = OUT_ROOT / _safe(case_name)

    def _compose_ic() -> Tuple[str, str]:
        title = f"IC Report — {case_name}"
        ic_map = ss.get("ic_map", {})

        raw_ten = ss.get("ten_steps") or {}
        scores = raw_ten.get("scores") or [5] * len(TEN_STEPS)
        narrs = raw_ten.get("narratives") or [f"{s}: tbd" for s in TEN_STEPS]
        ten = {"scores": scores, "narratives": narrs}

        b: List[str] = []
        interpreted = ss.get("combined_text", "").strip() or ss.get("narrative", "(no summary)")
        b.append(f"Executive Summary\n\n{interpreted}\n")
        if not PUBLIC_MODE:
            b.append(f"Evidence Quality: ~{ss.get('evidence_quality', 0)}% coverage (heuristic)\n")

        b.append("Four-Leaf Analysis")
        for leaf in ["Human", "Structural", "Customer", "Strategic Alliance"]:
            row = ic_map.get(leaf, {"tick": False, "narrative": "", "score": 0.0})
            tail = "" if PUBLIC_MODE else f" (score: {row.get('score', 0.0)})"
            b.append(f"- {leaf}: {'✓' if row.get('tick') else '•'} — {row.get('narrative', '')}{tail}")

        b.append("\nTen-Steps Readiness")
        for s, n in zip(TEN_STEPS, ten["narratives"]):
            b.append(f"- {n}")

        b.append("\nNotes")
        b.append(
            "This document is provided for high-level evaluation only."
            if PUBLIC_MODE
            else "CONFIDENTIAL. Advisory-first; company and LIP review required for final scoring, licensing design and accounting treatment."
        )
        return title, "\n".join(b)

    def _compose_lic() -> Tuple[str, str]:
        """
        Licensing-focused report, using the 6 SME-style questions plus IC map.
        Keeps language simple so it can be shared with SMEs / spin-outs if needed.
        """
        title = f"Licensing Readiness Report — {case_name}"

        size = ss.get("company_size", "Not specified")
        sector = ss.get("sector", "Other")

        # Pull the six company-page answers (re-using existing keys)
        why_service = (ss.get("why_service", "") or "").strip()        # "What are you hoping this help will achieve?"
        stage = (ss.get("stage", "") or "").strip()                    # "Where is the idea / product today?"
        plan_s = (ss.get("plan_s", "") or "").strip()                  # Short-term focus
        plan_m = (ss.get("plan_m", "") or "").strip()                  # Medium-term focus
        plan_l = (ss.get("plan_l", "") or "").strip()                  # Long-term focus
        markets_why = (ss.get("markets_why", "") or "").strip()        # Who would benefit / what markets
        sale_price_why = (ss.get("sale_price_why", "") or "").strip()  # What would feel “fair” if a partner asked tomorrow

        ic_map = ss.get("ic_map", {}) or {}
        evidence_quality = int(ss.get("evidence_quality", 0))

        # Simple flags for narrative
        structural_row = ic_map.get("Structural", {"tick": False, "score": 0.0})
        structural_ready = bool(structural_row.get("tick"))
        structural_score = float(structural_row.get("score", 0.0))

        customer_row = ic_map.get("Customer", {"tick": False, "score": 0.0})
        customer_ready = bool(customer_row.get("tick"))

        alliance_row = ic_map.get("Strategic Alliance", {"tick": False, "score": 0.0})
        alliance_ready = bool(alliance_row.get("tick"))

        b: List[str] = []

        # 1. Overview & purpose
        b.append(f"1. Overview & Purpose\n")
        b.append(
            f"{case_name} is described as a {size} organisation operating in {sector}. "
            "This report summarises how ready the company (or project/spin-out) is to license or share its know-how, "
            "technology, data or content, and what practical next steps are recommended."
        )
        if why_service:
            b.append(f"\nMain reason for using this service (in your own words):\n- {why_service}\n")
        else:
            b.append("\nMain reason for using this service has not been fully described yet.\n")

        # 2. Goals for licensing / partnering
        b.append("\n2. Goals for Licensing or Partnering\n")
        if any([plan_s, plan_m, plan_l]):
            if plan_s:
                b.append(f"- Next 0–6 months: {plan_s}")
            if plan_m:
                b.append(f"- Next 6–24 months: {plan_m}")
            if plan_l:
                b.append(f"- Beyond 24 months: {plan_l}")
        else:
            b.append(
                "- Short, medium and longer-term goals have not yet been set in detail. "
                "A short planning session is recommended to agree what ‘good’ looks like for the next 2–3 years."
            )

        # 3. Asset & readiness snapshot
        b.append("\n3. What is Being Licensed & How Ready It Is\n")
        if stage:
            b.append(f"- Current stage (how you described it): {stage}")
        else:
            b.append("- Current stage of the idea / product has not yet been clearly described.")

        if structural_ready:
            b.append(
                f"\nBased on the uploaded documents, there are signs of **codified assets** "
                f"(contracts, methods, processes, datasets, training, software, etc.) "
                f"that are suitable for licensing. Structural Capital scored around {structural_score:.1f} in the IC map, "
                "which suggests there is something real to license — not just an idea."
            )
        else:
            b.append(
                "\nThe IC analysis suggests that codified assets (contracts, processes, datasets, content, software) "
                "are not yet clearly evidenced. Before licensing, it will help to gather and tidy the key documents "
                "into a simple register (what exists, who owns it, and where it is stored)."
            )

        if customer_ready or alliance_ready:
            bullets = []
            if customer_ready:
                bullets.append("customer relationships or repeat business")
            if alliance_ready:
                bullets.append("partners, universities, suppliers or other strategic allies")
            b.append(
                "\nThere is also evidence of value sitting in "
                + " and ".join(bullets)
                + ", which can support partnership-based licensing and pilot deals."
            )

        if evidence_quality:
            b.append(
                f"\nEvidence quality from the file mix is estimated at about {evidence_quality}% "
                "(heuristic). More structured documents will increase confidence for investors, partners "
                "and auditors."
            )

        # 4. Who could benefit (markets & partner types)
        b.append("\n4. Who Could Benefit (Markets & Partner Types)\n")
        if markets_why:
            b.append(
                "In your own words, you see the best fit with:\n"
                f"- {markets_why}\n"
            )
            b.append(
                "Licensing professionals can translate this into a short list of **partner types** "
                "(for example: pilot customers, distributors, technology partners, universities, NGOs or government programmes)."
            )
        else:
            b.append(
                "- Target markets and partner types are not yet described. A quick exercise to map 3–5 ideal partner profiles "
                "will make licensing discussions much easier."
            )

        # 5. Pricing & “fair” outcome
        b.append("\n5. What Would Feel ‘Fair’ if a Partner Asked Tomorrow?\n")
        if sale_price_why:
            b.append(
                "You described a fair outcome as:\n"
                f"- {sale_price_why}\n"
            )
            b.append(
                "This is a good starting point. A licensing professional can now work backwards from this to design:\n"
                "- simple, transparent fee or royalty structures; and\n"
                "- options for pilots or community/low-fee access where that makes sense."
            )
        else:
            b.append(
                "- A clear view of what would feel ‘fair’ has not been written down yet. "
                "A short internal discussion (what minimum would we accept, what would feel like a win, "
                "and what are we *not* willing to do) will speed up negotiations later."
            )

        # 6. Suggested licensing directions (plain language)
        b.append("\n6. Suggested Licensing Directions (Plain Language)\n")
        b.append(
            "Based on the information so far, the following **families of licensing options** are likely to be relevant. "
            "Exact terms will depend on the asset, partners and country-specific law."
        )
        b.append(
            "- **Paid usage licences** – a partner pays a fee or royalty to use your software, content, method, data or brand "
            "under agreed conditions (field of use, territory, time period)."
        )
        b.append(
            "- **Pilot or project licences** – a short, time-bound licence to test the asset in a real setting, usually with a "
            "limited number of users or sites. Useful for proofs-of-concept and early adoption."
        )
        b.append(
            "- **Co-creation / joint development** – you and one or more partners build something together. Ownership of the "
            "new results is shared and clearly written down (who owns what, how revenue or savings are shared)."
        )
        b.append(
            "- **Access or community licences** – low-fee or free access for specific groups (for example schools, farmers, "
            "public bodies or communities) where the main value is impact, data, learning or reputation."
        )
        b.append(
            "- **Data / algorithm access** – controlled use of your datasets, indices, risk scores or algorithms, with clear rules "
            "on how they are used, who can see what, and how results are reported back."
        )

        b.append(
            "\nA licensing professional can combine these building blocks into a small number of realistic options that fit "
            "your situation (for example: one commercial licence, one pilot licence and one access-for-impact licence)."
        )

        # 7. Immediate next steps
        b.append("\n7. Immediate Next Steps\n")
        b.append(
            "- Create a simple **asset list**: what you want to license (software, content, training, methods, datasets, indices, "
            "brand elements, etc.), where each item is stored and who currently owns/controls it."
        )
        b.append(
            "- Identify 3–5 **ideal partner types** and 3–5 concrete organisations to approach first."
        )
        b.append(
            "- Decide your **red lines and must-haves**: what you will not give away, how you want your work to be credited, and "
            "which partners or sectors are off-limits."
        )
        b.append(
            "- Work with a Licensing & Intangibles Partner (LIP) and legal adviser to turn this report into draft licence terms "
            "and a small set of template agreements."
        )

        # 8. Status & disclaimer
        b.append("\n8. Status & Disclaimer\n")
        b.append(
            "This report is an advisory draft only. It does **not** replace legal advice. All licensing decisions and legal "
            "documents must be reviewed and approved by qualified legal counsel in the relevant country or countries."
        )

        return title, "\n".join(b)

    c1, c2 = st.columns(2)
    with c1:
        if st.button("Generate IC Report (DOCX/TXT)", key="btn_ic"):
            title, body = _compose_ic()
            data, fname, mime = _export_bytes(title, body)
            path, msg = _save_bytes(case_folder, fname, data)
            st.download_button(
                "⬇️ Download IC Report",
                data,
                file_name=fname,
                mime=mime,
                key="dl_ic",
            )
            (st.success if path else st.warning)(msg)
    with c2:
        if st.button("Generate Licensing Report (DOCX/TXT)", key="btn_lic"):
            title, body = _compose_lic()
            data, fname, mime = _export_bytes(title, body)
            path, msg = _save_bytes(case_folder, fname, data)
            st.download_button(
                "⬇️ Download Licensing Report",
                data,
                file_name=fname,
                mime=mime,
                key="dl_lic",
            )
            (st.success if path else st.warning)(msg)

    st.caption("Server save root: disabled (public mode)" if PUBLIC_MODE else f"Server save root: {OUT_ROOT}")

# 6) LICENSING TEMPLATES
elif page == "Licensing Templates":
    st.header("Licensing Templates (editable DOCX/TXT)")
    case = ss.get("case_name", "Untitled Company")
    sector = ss.get("sector", "Other")

    template = st.selectbox(
        "Choose a template:",
        ["FRAND Standard", "Co-creation (Joint Development)", "Knowledge (Non-traditional)"],
        index=0,
    )

    if st.button("Generate template", key="btn_make_template"):
        disclaimer = (
            "STATUS: Non-binding draft template. This document is provided for internal discussion only and must be "
            "reviewed and adapted by qualified legal counsel before signature or implementation.\n\n"
        )

        if template == "FRAND Standard":
            title = f"FRAND Standard template — {case}"
            body = (
                f"FRAND Standard Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Purpose & Scope\n"
                "   - Define the licensed technology, dataset, index, software or method.\n"
                "   - Specify fields of use, territories and permitted users.\n\n"
                "2. Access & FRAND-Informed Principles\n"
                "   - Describe how fees (including possible royalty-free tiers) are set on a fair, reasonable and\n"
                "     non-discriminatory basis across comparable users.\n"
                "   - Include access considerations for social, community or public-good partners where relevant.\n\n"
                "3. Financial Terms\n"
                "   - Royalty or fee structure (e.g. per unit, per user, revenue share, or capped fees).\n"
                "   - Invoicing, payment schedule, late payment and currency terms.\n\n"
                "4. Governance, Reporting & Audit\n"
                "   - Licensee reporting obligations (KPIs, usage, sublicensing, ESG/stakeholder indicators if applicable).\n"
                "   - Audit rights and frequency; treatment of under-reporting or non-compliance.\n\n"
                "5. IP Ownership & Improvements\n"
                "   - Background IP ownership and reservation of rights.\n"
                "   - Treatment of improvements, derivative works and feedback.\n\n"
                "6. Compliance, Term & Termination\n"
                "   - Conditions for suspension or termination (breach, non-payment, misuse).\n"
                "   - Survival of key clauses (confidentiality, audit, data protection).\n\n"
                "7. Data, AI & Regulatory Considerations (if applicable)\n"
                "   - Data protection, confidentiality and AI/automated decision-making safeguards.\n"
                "   - Reference to applicable laws, standards or regulatory guidance.\n\n"
                "8. Governing Law & Dispute Resolution\n"
                "   - Governing law (e.g. EU Member State law).\n"
                "   - Mechanism for dispute resolution (negotiation, mediation, arbitration, courts).\n"
            )
        elif template == "Co-creation (Joint Development)":
            title = f"Co-creation template — {case}"
            body = (
                f"Co-creation / Joint Development Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Parties, Purpose & Project Definition\n"
                "   - Identify all parties and describe the joint development project and objectives.\n\n"
                "2. Background IP\n"
                "   - List key Background IP contributed by each party and conditions of use.\n\n"
                "3. Foreground IP & Ownership Structure\n"
                "   - Define Foreground IP and allocate ownership shares (e.g. 50/50 or by contribution).\n"
                "   - Set rules for registration, maintenance and enforcement of Foreground IP.\n\n"
                "4. Contributions, Resources & Cost Sharing\n"
                "   - Describe personnel, facilities, data and funding contributed by each party.\n"
                "   - Agree cost-sharing mechanisms for development and exploitation.\n\n"
                "5. Commercialisation & Revenue Sharing\n"
                "   - Outline commercialisation routes (direct sales, licensing, joint ventures).\n"
                "   - Define revenue sharing mechanisms, including treatment of royalty-free access for selected stakeholders.\n\n"
                "6. Publication, Academic Use & Confidentiality\n"
                "   - Academic and scientific publication rights (timing, review, attribution).\n"
                "   - Confidentiality obligations and carve-outs.\n\n"
                "7. Governance, Decision-Making & Dispute Resolution\n"
                "   - Governance structure (steering committee, decision rules).\n"
                "   - Escalation and dispute resolution framework.\n\n"
                "8. Term, Exit & Transition\n"
                "   - Term and conditions for early termination.\n"
                "   - Exit options (buy-out, assignment, break clauses) and treatment of Foreground IP on exit.\n"
            )
        else:
            title = f"Knowledge licence (non-traditional) — {case}"
            body = (
                f"Knowledge Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Knowledge Asset Definition\n"
                "   - Describe the codified know-how (e.g. methods, training materials, playbooks, indices.   m , checklists).\n\n"
                "2. Scope of Licence & Field of Use\n"
                "   - Specify permitted uses (internal training, consulting, product development, policy design, etc.).\n"
                "   - Clarify any restricted uses and prohibited activities.\n\n"
                "3. Access & Stakeholder Considerations\n"
                "   - Define standard and, where appropriate, community or social-benefit access tiers.\n"
                "   - Reference a FRAND-informed approach to fairness, reasonableness and non-discrimination\n"
                "     across employees, investors, customers, partners, suppliers, communities and nature.\n\n"
                "4. Attribution & Moral Rights\n"
                "   - Conditions for attribution (branding, acknowledgements, citation requirements).\n\n"
                "5. Confidentiality, Data Protection & Safeguards\n"
                "   - Treatment of confidential information and personal data.\n"
                "   - Safeguards where AI or automated decision-making is involved.\n\n"
                "6. Term, Revocation & Review\n"
                "   - Duration, renewal and review points.\n"
                "   - Conditions for revocation or modification of the licence.\n\n"
                "7. Governing Law & Dispute Resolution\n"
                "   - Governing law.\n"
                "   - Mechanism for addressing disputes.\n"
            )

        data, fname, mime = _export_bytes(title, body)
        folder = OUT_ROOT / _safe(case)
        path, msg = _save_bytes(folder, fname, data)
        st.download_button(
            "⬇️ Download Template",
            data,
            file_name=fname,
            mime=mime,
            key="dl_tpl",
        )
        (st.success if path else st.warning)(msg)

# 7) LIP ASSISTANT (beta)
elif page == "LIP Assistant":
    st.header("LIP Assistant (beta)")
    st.caption(
        "A lightweight helper for the Licensing & Intangibles Partner. "
        "It re-uses the IC narrative and licensing logic locally — no external AI calls are made in this demo."
    )

    if ss.get("combined_text", ""):
        st.success("IC narrative available from Analyse Evidence / LIP Console.")
    else:
        st.warning("No IC narrative stored yet. Run **Analyse Evidence** first for best results.")

    context_choice = st.selectbox(
        "Which context should the LIP Assistant focus on?",
        ["IC findings", "Licensing options", "Both"],
        index=2,
    )

    question = st.text_area(
        "Your question",
        "",
        height=120,
        help="Example: 'Which assets look IAS 38-ready?' or 'How could we structure a royalty-free community licence?'",
    )

    if st.button("Ask LIP Assistant"):
        if not question.strip():
            st.error("Please enter a question.")
        else:
            ic_text = ss.get("combined_text", "")
            # Short licensing explainer from the report logic
            lic_snippet = (
                "The tool supports revenue licences, access/community licences, co-creation/joint development, "
                "defensive/cross-licence arrangements and data/algorithm licences, all treated through a FRAND-informed "
                "lens across the Seven Stakeholder Model."
            )

            answer_parts: List[str] = []
            answer_parts.append(f"**Question received:** {question.strip()}")

            q_lower = question.lower()

            if context_choice in ("IC findings", "Both"):
                answer_parts.append("**IC & Structural Capital view (IAS 38):**")
                if "structural" in q_lower or "ias 38" in q_lower:
                    answer_parts.append(
                        "- Structural Capital is treated as the home for explicit, documented assets "
                        "(contracts, SOPs, registers, CRM, datasets, board packs). Where both tacit and explicit "
                        "signals exist, Structural is deliberately weighted to dominate, reflecting IAS 38 audit-readiness."
                    )
                if ic_text:
                    answer_parts.append(
                        "- Key IC narrative extract:\n\n"
                        + ic_text[:600]
                        + ("..." if len(ic_text) > 600 else "")
                    )
                else:
                    answer_parts.append(
                        "- No stored IC narrative yet. Run the analysis and revisit this view to see a richer explanation."
                    )

            if context_choice in ("Licensing options", "Both"):
                answer_parts.append("**Licensing & FRAND view:**")
                if "frand" in q_lower:
                    answer_parts.append(
                        "- In this tool, FRAND is a *design principle* rather than a legal certification. "
                        "It guides how pricing, access tiers and governance are framed across different stakeholders, "
                        "including royalty-free and community licences."
                    )
                answer_parts.append("- " + lic_snippet)

            answer_parts.append(
                "_Next step: a Licensing & Intangibles Partner can tailor these insights into concrete clauses using the "
                "reports and templates generated in the other pages._"
            )

            ss["lip_history"].append({"q": question.strip(), "a": "\n\n".join(answer_parts)})
            st.markdown("\n\n".join(answer_parts))

    if ss.get("lip_history"):
        with st.expander("Previous LIP Assistant exchanges"):
            for i, entry in enumerate(reversed(ss["lip_history"]), start=1):
                st.markdown(f"**Q{i}:** {entry['q']}")
                st.markdown(entry["a"])
                st.markdown("---")

## 8) GLOSSARY — simple static page for key terms
elif page == "Glossary":
    st.header("Glossary – key terms & acronyms")

    st.markdown(
        """
### 1. Intellectual Capital, Intangible Assets and Areopa models

**Intellectual Capital (IC)**  
The total “know-how value” of an organisation – people, structures, relationships and partnerships that create value but do not sit neatly on the balance sheet. In this tool, IC is analysed using the 4-Leaf Model and Ten-Steps.

**Intangible Assets (IA)**  
The explicit, documented parts of IC that can potentially be recognised as assets under accounting rules (for example software, licences, brands, contracts, datasets). This tool helps identify which parts of IC are close to being IA.

**4-Leaf Model**  
Areopa’s IC framework. It groups Intellectual Capital into four “leaves”:  
- **Human Capital** – people, skills, experience and culture.  
- **Structural Capital** – documented processes, systems, IP, contracts and data.  
- **Customer Capital** – relationships with customers, users and channels.  
- **Strategic Alliance Capital** – value created with partners, universities, suppliers, funders and other allies.

**Human Capital (HC)**  
The knowledge, experience, skills and culture held by people. Often tacit (in heads and habits). It becomes more valuable – and more auditable – when training, competency frameworks and safety records are documented.

**Structural Capital (SC)**  
The “codified” side of the business: contracts, licences, SOPs, manuals, datasets, software, board packs, registers, QMS, CRM and other repeatable systems. In this tool, Structural Capital is treated as the main home for potential IAS 38-ready intangible assets.

**Customer Capital (CC)**  
The strength and quality of customer relationships: who you sell to, how often they return, what contracts you have, and how you manage the pipeline and renewals (often through CRM or account plans).

**Strategic Alliance Capital (SAC)**  
Value built with partners and allies: universities, research institutes, innovation hubs, corporates, funders, governments, NGOs, suppliers and integrators. MoUs, JV agreements and framework contracts are key artefacts here.

**Ten-Steps (IC lifecycle)**  
Areopa’s 10-step lifecycle for building and managing IC and IA. In this tool each step is given a readiness score (1–10):  
1. **Identify** – spot the knowledge, methods, software, data and relationships that create value.  
2. **Separate** – distinguish what is really new or special from what is generic or commodity.  
3. **Protect** – put in place IP, contracts, confidentiality and access controls where needed.  
4. **Safeguard** – make sure assets cannot be easily lost, copied or misused (back-ups, roles, policies).  
5. **Manage** – give someone responsibility to maintain and use the asset, not just “store” it.  
6. **Control** – set clear rights and permissions: who can use, change, license or sell the asset.  
7. **Use** – actually deploy the asset in products, services, partnerships or internal systems.  
8. **Monitor** – track performance, usage, risks and outcomes (KPIs, dashboards, board packs).  
9. **Value** – estimate financial value (for deals, investment, or accounting) using a suitable method.  
10. **Report** – explain the asset and its value clearly to boards, investors, auditors and regulators.

**IC Map**  
The combined picture of the 4-Leaf Model and Ten-Steps scores generated by the tool. It shows where IC is strong, where it is under-evidenced, and which assets are closest to being licensing-ready or IAS 38-ready.

---

### 2. Accounting, valuation and evidence concepts

**IAS 38 – International Accounting Standard 38 (Intangible Assets)**  
The main IFRS standard covering recognition and measurement of many intangible assets (such as software, brands, some IP and development costs). This tool does **not** give an accounting opinion, but it highlights Structural Capital that looks close to IAS 38 criteria (identifiable, controlled and expected to generate future economic benefits).

**Fair value (approach)**  
A valuation approach that estimates what a knowledgeable, willing buyer would pay in an arm’s-length transaction. In the tool, fair value thinking is used when discussing how explicit assets might support deals, licensing or investment – not just what they cost to create.

**Cost approach**  
A valuation method that focuses on how much it cost (or would cost) to build or replace an asset. Often underestimates value for high-impact IC where licensing, market access or unique positioning matter. The tool leans towards fair value logic rather than pure cost.

**Evidence quality (heuristic)**  
A percentage score used inside the app to summarise how strong the uploaded evidence set looks. It blends:  
- diversity of file types,  
- spread across the four IC leaves, and  
- how “codified” the artefacts appear to be.  
It is a guidance signal only – not an audit opinion.

**IA Register (Intangible Asset Register)**  
A structured list of key intangible assets: what they are, where they are recorded, who owns them, and how they are used. The tool encourages the creation of such a register from contracts, software, datasets, methods and brand assets that show up in Structural Capital.

**Structural Capital dominance**  
A design choice in this tool: when both tacit and explicit signals exist for an asset, Structural Capital is deliberately weighted higher. This reflects the practical reality that auditors and investors place more trust in clear, documented artefacts.

---

### 3. Roles, pages and tool components

**IC-LicAI**  
The name of this tool – an IC and Licensing Assistant for identifying, structuring and preparing intangible assets and licensing options, with an emphasis on Structural Capital and FRAND-aware licensing.

**VM – Value Manager**  
The person preparing evidence, running diagnostics and curating assumptions into IC and licensing reports. Often works closely with founders, TTOs or programme teams.

**LIP – Licensing & Intangibles Partner**  
The senior person (internal or external) who interprets results, designs licensing options and works with legal counsel. The LIP uses the LIP Console, Reports and Templates pages.

**Company / project page**  
The starting page where basic context is captured in SME-friendly language: what is being built, stage of development, who will use it, what already exists and how the team hopes to earn and share value.

**Analyse Evidence (Evidence Dashboard)**  
The page that pulls together uploaded files, runs the IC analysis, and shows radar charts and Ten-Steps readiness. It also surfaces the “evidence quality” score.

**Asset Verification**  
A human review page where the VM or LIP checks ownership, dates, validity and ESG claims for each document, and records follow-up notes. It does **not** replace legal review.

**LIP Console**  
A workspace for the LIP to edit the overall narrative, review the IC Map, and see key context from the company page before generating reports or templates.

**Reports page**  
Generates two main exports:  
- an **IC Report** – focused on IC, Ten-Steps and evidence coverage; and  
- a **Licensing Readiness Report** – aimed at SMEs, summarising what can be licensed and practical next steps.

**Licensing Templates page**  
Creates editable draft templates (DOCX/TXT) for three families of agreements: FRAND-style standard licences, co-creation/joint development arrangements, and non-traditional knowledge licences.

**LIP Assistant (beta)**  
A lightweight Q&A helper that re-uses the IC narrative and licensing logic locally (no external AI calls in this demo). It is designed to assist the LIP in thinking through options, not to draft final legal text.

---

### 4. Licensing concepts used in the tool

**Licensing / Licence**  
A legal permission to use software, data, content, methods, brands or other assets under defined conditions (for example where, for how long, and for which users). Licensing can generate revenue, support pilots, or enable wider access.

**Field of use**  
The specific application or sector where a licence applies (for example “healthcare research only”, “agricultural advisory services” or “educational use”). Narrow fields of use allow different licences in different markets.

**Territory**  
The geographic scope of a licence (for example one country, a region such as EU or East Africa, or worldwide).

**FRAND – Fair, Reasonable and Non-Discriminatory**  
In this tool, FRAND is used as a design principle, not a formal legal status. It guides how pricing, access tiers and governance are framed so that similar users are treated fairly, and community or public-good access can be justified and described.

**Standard licence (FRAND-informed)**  
A more traditional commercial licence with transparent prices or royalty terms, shaped by FRAND ideas and, where relevant, by stakeholder and ESG considerations.

**Co-creation / Joint Development**  
A partnership where two or more parties jointly develop new assets. Background IP (what each party brings) is kept separate from Foreground IP (results created together). Revenue sharing and ownership rules are explicitly written down.

**Knowledge licence (non-traditional)**  
A licence focused on codified know-how – methods, training content, playbooks, indices, checklists – rather than classic patents or software only. Can be structured for internal use, consulting, policy work or education.

**Pilot licence**  
A short-term or limited-scope licence used to test a product, service or method with early adopters. Often restricted by time, number of users or locations.

**Access / community licence**  
A licence that provides low-cost or free access for specific groups (for example schools, farmers, public bodies or community organisations) where the main value is impact, learning, data or reputation rather than direct revenue.

**Data / algorithm licence**  
A licence that allows controlled use of datasets, indices, risk scores or algorithms (for example via APIs or dashboards) with clear rules on privacy, attribution, reuse and feedback of results.

**Background IP**  
Existing intellectual property that each party brings into a project before any new work is done (for example an existing platform, method, dataset or brand).

**Foreground IP**  
New IP created during a project or collaboration. A co-creation or joint development agreement should clearly define who owns it, how it is licensed and how revenues or savings are shared.

**Royalty**  
A payment calculated as a percentage of revenue, profit or usage (for example per unit or per user) that the licensee pays to the licensor under a licensing agreement.

**Revenue share**  
A broader concept than royalty – an agreed split of income (sometimes including grants or savings) between parties to a licence or partnership.

**NDA – Non-Disclosure Agreement**  
A confidentiality agreement that controls how sensitive information may be shared and used, usually before or alongside licensing or co-creation discussions.

**MoU – Memorandum of Understanding**  
A non-binding (or partly binding) framework that sets out intentions and main principles for collaboration. Often a stepping stone before a full contract.

**Joint Venture (JV)**  
A separate structure (often a new company) owned by two or more parties to exploit specific assets or markets, with detailed agreements on contributions, governance and profit sharing.

---

### 5. ESG, stakeholders and policy language

**ESG – Environmental, Social and Governance**  
A set of factors used by investors, regulators and stakeholders to assess a company’s wider impact. In this tool, ESG language is linked to concrete assets (for example policies, datasets, procedures and dashboards) rather than slogans.

**Seven Stakeholder Model (SSM)**  
A framework (Sugai–Weir) for mapping value creation and impact across seven stakeholder groups:  
1. Employees and workers  
2. Investors and lenders  
3. Customers and users  
4. Partners, suppliers and intermediaries  
5. Communities and public bodies  
6. The natural environment / nature  
7. Future generations (often reflected through long-term environmental and social impacts)

**Double materiality**  
The idea that we must look both at:  
- how sustainability issues affect the company’s financial performance; and  
- how the company’s activities affect people, communities and the environment.  
The tool supports this by tying ESG and stakeholder claims back to specific intangible assets and evidence.

**Impact or community access licences**  
Licences that deliberately build in reduced fees, open access or shared ownership for certain stakeholder groups (for example farmers, schools or low-income communities) while still protecting core IP and commercial value.

---

### 6. Project and funding references

**IMPAC3T-IP**  
The Horizon Europe project under which this tool has been developed. It focuses on practical methods for valuing Intellectual Capital, ESG and licensing in SMEs and ecosystems.

**Horizon Europe**  
The European Union’s main research and innovation funding programme. IMPAC3T-IP is funded under Grant Agreement No. 101135832.

If a term appears in the tool and is **not** in this glossary, it can be added or clarified in future versions.
"""
    )
