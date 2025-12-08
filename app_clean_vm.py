# app_clean.py — IC-LicAI Expert Console (Structural + FRAND + LIP Assistant)
# Adds: DOCX/PPTX extraction, weighted IC signal engine, interpreted narrative,
# radar dashboard, CSV semantic extraction, robust company-context auto-split,
# IAS 38 Structural Capital emphasis, FRAND-aware licensing templates,
# Seven Stakeholder / ESG narrative, LIP Console, and LIP Assistant (beta).

from __future__ import annotations
import io, os, tempfile, re, csv
from pathlib import Path
from typing import Dict, Any, List, Tuple, Optional
from dataclasses import dataclass

import streamlit as st
import plotly.graph_objects as go  # for radar charts
import requests


@dataclass
class VMAssumption:
    key: str
    label: str
    narrative: str
    rationale: str
    category: str
    source_signals: List[str]
    confidence: str
    include: bool = True

HAVE_PDF = False
try:
    from PyPDF2 import PdfReader  # type: ignore
    HAVE_PDF = True
except Exception:
    HAVE_PDF = False

# -------------------- PROJECT LOGOS -----------------
BASE_DIR = Path(__file__).parent  # folder where app_clean_vm.py lives

IMPACT3T_LOGO_PATH = BASE_DIR / "demo_assets" / "impact3t_logo.png"
EU_FLAG_PATH = BASE_DIR / "demo_assets" / "eu_flag.png"

# -------------------- MODE / AUTH --------------------
# VM version: always PRIVATE (internal, passphrase required)
APP_MODE = "PRIVATE"

PUBLIC_MODE: bool = False
REQUIRE_PASS: bool = True

# ---------------- DOCX/PPTX/PDF optional ----------------
HAVE_DOCX = False
HAVE_PPTX = False

try:
    from docx import Document  # type: ignore
    HAVE_DOCX = True
except Exception:
    HAVE_DOCX = False

try:
    from pptx import Presentation  # type: ignore
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False
    
HAVE_PDF = False
try:
    from PyPDF2 import PdfReader  # type: ignore
    HAVE_PDF = True
except Exception:
    HAVE_PDF = False
    
# ------------------ THEME ----------------------------
# IMPAC3T-IP inspired palette (no yellow / gold)
PRIMARY_NAVY   = "#003B70"  # deep blue, EU-friendly
ACCENT_BLUE    = "#00A7E1"  # bright accent from IMPAC3T-style slides
ACCENT_PURPLE  = "#6E4BAF"  # secondary accent if needed later
BG_MAIN        = "#F5F7FB"  # soft neutral background
SIDEBAR_BG     = "#0F2F56"  # dark blue sidebar
CARD_BG        = "#FFFFFF"  # white cards / blocks
TEXT_MAIN      = "#1F2933"  # dark grey text

st.set_page_config(page_title="IC-LicAI Expert Console", layout="wide")

st.markdown(
    f"""
<style>
  /* Overall app background */
  .stApp {{
    background: {BG_MAIN};
  }}

  /* Main content container */
  .block-container {{
    max-width: 1250px;
    padding-top: 1.2rem;
    padding-bottom: 2rem;
  }}

  /* Title bar at the top */
  .ic-title-bar {{
    background: {PRIMARY_NAVY};
    color: #FFFFFF;
    font-weight: 800;
    font-size: 32px;
    padding: 16px 22px;
    border-radius: 10px;
    letter-spacing: 0.2px;
    margin: 10px 0 20px 0;
    box-shadow: 0 2px 6px rgba(0,0,0,0.08);
  }}

  /* Primary buttons */
  .stButton>button {{
    background: {ACCENT_BLUE} !important;
    color: #FFFFFF !important;
    border-radius: 8px !important;
    border: 0 !important;
    padding: 0.55rem 1rem !important;
    font-weight: 700 !important;
  }}

  .stButton>button:hover {{
    background: {PRIMARY_NAVY} !important;
  }}

  /* Sidebar styling */
  section[data-testid="stSidebar"] {{
    background: {SIDEBAR_BG};
  }}

  section[data-testid="stSidebar"] h1,
  section[data-testid="stSidebar"] h2,
  section[data-testid="stSidebar"] h3,
  section[data-testid="stSidebar"] p,
  section[data-testid="stSidebar"] label,
  section[data-testid="stSidebar"] span {{
    color: #E7F0FF !important;
  }}

  .stRadio div[role="radiogroup"] label {{
    color: #E7F0FF !important;
  }}

  /* Make cards / expanders look clean on light background */
  .stExpander, .stTabs [data-baseweb="tab"] {{
    background: {CARD_BG};
  }}
</style>
""",
    unsafe_allow_html=True,
)

st.markdown(
    '<div class="ic-title-bar">IC-LicAI Expert Console</div>',
    unsafe_allow_html=True,
)
st.caption("INTERNAL VERSION — FOR REAL EVIDENCE (PASS-PHRASE PROTECTED)")
# ------------------ AUTH GATE ------------------------
def _auth_gate() -> None:
    if not REQUIRE_PASS:
        return
    secret = st.secrets.get("APP_KEY", None) or os.environ.get("APP_KEY", None)
    if not secret:
        with st.expander("Access control"):
            st.info("Optional passphrase: set st.secrets['APP_KEY'] or env APP_KEY.")
        return
    key = st.text_input("Enter access passphrase", type="password")
    if not key:
        st.stop()
    if key != secret:
        st.error("Incorrect passphrase.")
        st.stop()

_auth_gate()

# --------------- WRITABLE ROOT -----------------------
def _detect_writable_root() -> Path:
    for p in [Path("./out"), Path(os.path.expanduser("~")) / "out", Path(tempfile.gettempdir()) / "ic-licai-out"]:
        try:
            p.mkdir(parents=True, exist_ok=True)
            t = p / ".touch"
            t.write_text("ok", encoding="utf-8")
            t.unlink()
            return p
        except Exception:
            continue
    return Path(tempfile.gettempdir())

OUT_ROOT = _detect_writable_root()


def _ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)


def _safe(name: str) -> str:
    return "".join(c for c in (name or "").strip() if c.isalnum() or c in (" ", "_", "-", ".")).strip().replace(" ", "_")


def _export_bytes(title: str, body: str) -> Tuple[bytes, str, str]:
    base = _safe(title) or "ICLicAI_Report"
    if HAVE_DOCX:
        doc = Document()
        if not PUBLIC_MODE:
            doc.add_paragraph().add_run("CONFIDENTIAL — Internal Evaluation Draft (No Distribution)").bold = True
        doc.add_heading(title, 0)
        for para in body.split("\n\n"):
            doc.add_paragraph(para)
        bio = io.BytesIO()
        doc.save(bio)
        return (
            bio.getvalue(),
            f"{base}.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    if not PUBLIC_MODE:
        body = "CONFIDENTIAL — Internal Evaluation Draft (No Distribution)\n\n" + body
    return body.encode("utf-8"), f"{base}.txt", "text/plain"

def _save_bytes(folder: Path, name: str, data: bytes) -> Tuple[Optional[Path], str]:
    if PUBLIC_MODE:
        return None, "Public mode: server save disabled (download only)."
    try:
        _ensure_dir(folder)
        p = folder / name
        p.write_bytes(data)
        return p, f"Saved to {p}"
    except Exception as e:
        return None, f"Server save skipped ({type(e).__name__}: {e}). Download only."
        
# --------------- EVIDENCE EXTRACTION -----------------
TEXT_EXT = {".txt"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
CSV_EXT  = {".csv"}
PDF_EXT  = {".pdf"}  # filename cue for PDFs

def _extract_text_docx(data: bytes) -> str:
    if not HAVE_DOCX:
        return ""
    try:
        bio = io.BytesIO(data)
        doc = Document(bio)
        parts: List[str] = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if txt:
                parts.append(txt)
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                line = " | ".join((cell.text or "").strip() for cell in row.cells)
                if line.strip():
                    parts.append(line)
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_text_pptx(data: bytes) -> str:
    if not HAVE_PPTX:
        return ""
    try:
        bio = io.BytesIO(data)
        prs = Presentation(bio)
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = (shape.text or "").strip()
                    if txt:
                        parts.append(txt)
            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                nt = (slide.notes_slide.notes_text_frame.text or "").strip()
                if nt:
                    parts.append(nt)
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_text_pdf(data: bytes) -> str:
    """
    Extract basic text from a PDF using PyPDF2.
    (Scanned/image-only PDFs will still come back empty.)
    """
    if not HAVE_PDF:
        return ""
    try:
        bio = io.BytesIO(data)
        reader = PdfReader(bio)
        parts: List[str] = []
        for page in reader.pages:
            txt = page.extract_text() or ""
            txt = txt.strip()
            if txt:
                parts.append(txt)
        return "\n".join(parts)
    except Exception:
        return ""


def _pdf_review_hints(data: bytes, name: str) -> List[str]:
    """
    Light-weight scan of a PDF to suggest pages a Value Manager should review.
    We look for pages that mention tables/figures, IP, contracts, KPIs, markets, etc.,
    and return human-readable hints with page numbers.

    This does NOT change scoring – it is just guidance.
    """
    hints: List[str] = []

    if not HAVE_PDF:
        return hints

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception:
        return hints

    total_pages = len(reader.pages)

    # Simple keyword buckets – tweakable later.
    TABLE_KEYS = ["table", "figure", "diagram", "chart", "exhibit", "kpi", "metric"]
    PROCESS_KEYS = ["process", "workflow", "procedure", "protocol", "sop", "ipr process", "governance"]
    IP_KEYS = [
        "intellectual property",
        "ip register",
        "ipr",
        "patent",
        "trademark",
        "trade mark",
        "copyright",
        "licence",
        "license",
        "licensing",
        "contract",
        "agreement",
        "mou",
    ]
    SALES_KEYS = ["revenue", "turnover", "sales", "pipeline", "order book", "customer contract"]
    MARKET_KEYS = [
        "market",
        "segment",
        "customer",
        "client",
        "region",
        "country",
        "go-to-market",
        "g2m",
        "competition",
        "competitor",
    ]
    TECH_KEYS = [
        "technology",
        "platform",
        "software",
        "saas",
        "ai",
        "algorithm",
        "model",
        "index",
        "indices",
        "data platform",
    ]

    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").lower()
        except Exception:
            continue

        if not text.strip():
            continue

        categories: List[str] = []

        if any(k in text for k in TABLE_KEYS):
            categories.append("tables / figures / KPIs")
        if any(k in text for k in PROCESS_KEYS):
            categories.append("process / workflow / governance")
        if any(k in text for k in IP_KEYS):
            categories.append("IP / contracts / IA register")
        if any(k in text for k in SALES_KEYS):
            categories.append("revenue / sales / contracts")
        if any(k in text for k in MARKET_KEYS):
            categories.append("markets / customers / competitors")
        if any(k in text for k in TECH_KEYS):
            categories.append("technology / platform / indices")

        if categories:
            cat_txt = ", ".join(categories)
            hints.append(
                f"Page {idx}: check for {cat_txt} – see if this is explicit Structural Capital or still tacit."
            )

    # If nothing specific was found, still give one generic hint.
    if not hints and total_pages > 0:
        hints.append(
            "No specific keyword pages detected – skim the PDF for any contracts, KPIs, registers, "
            "or process diagrams that might indicate Structural Capital."
        )

    return hints


def _extract_text_csv(raw: bytes, name: str) -> str:
    """
    CSV semantic extraction: surface headers + a few rows so SME/ESG words
    like 'safety', 'policy', 'contract', 'governance', 'emissions', etc.
    are visible to the heuristics.
    """
    try:
        decoded = raw.decode("utf-8", errors="ignore")
        rows = list(csv.reader(decoded.splitlines()))
        if not rows:
            return ""
        headers = [h.strip() for h in rows[0] if h.strip()]
        cells: List[str] = []
        for row in rows[1:11]:
            cells.extend([c.strip() for c in row if c.strip()])
        header_txt = ", ".join(headers)
        cells_txt = "; ".join(cells)
        return f"CSV:{name}\nHeaders: {header_txt}\nRows: {cells_txt}"
    except Exception:
        try:
            return raw.decode("utf-8", errors="ignore")
        except Exception:
            return ""


def _read_text(files: List[Any]) -> Tuple[str, Dict[str, int], Dict[str, float]]:
    """
    Returns (combined_text, counts_by_ext, weights_used)
    Weights depend on artefact type (contract/JV > SOP/KMP > specs/slides > culture).
    """
    chunks: List[str] = []
    counts: Dict[str, int] = {}
    weights_used: Dict[str, float] = {}

    NAME_WEIGHTS: List[Tuple[str, float]] = [
        ("contract", 1.0),
        ("msa", 1.0),
        ("sow", 0.9),
        ("sla", 0.9),
        ("agreement", 0.9),
        ("joint venture", 1.0),
        ("joint_venture", 1.0),
        ("jv", 1.0),
        ("mou", 1.0),
        ("grant", 0.9),
        ("licence", 0.9),
        ("license", 0.9),
        ("register", 0.9),
        ("knowledge_management", 0.8),
        ("kmp", 0.8),
        ("sop", 0.8),
        ("process", 0.8),
        ("safety", 0.8),
        ("protocol", 0.8),
        ("spec", 0.6),
        ("canvas", 0.6),
        ("bmc", 0.6),
        ("slides", 0.6),
        ("deck", 0.6),
        ("board_pack", 0.8),
        ("board", 0.8),
        ("pricing", 0.7),
        ("tariff", 0.7),
        ("dataset", 0.7),
        ("culture", 0.4),
        ("award", 0.4),
    ]
    EXT_DEFAULTS: Dict[str, float] = {
        ".docx": 0.7,
        ".pptx": 0.6,
        ".txt": 0.5,
        ".csv": 0.6,
        ".pdf": 0.4,
    }

    # Reset PDF hints each run
    if "pdf_hints" not in st.session_state:
        st.session_state["pdf_hints"] = {}
    else:
        st.session_state["pdf_hints"] = {}

    for f in files or []:
        name = getattr(f, "name", "file")
        lower_name = str(name).lower()
        ext = Path(lower_name).suffix or "none"
        counts[ext] = counts.get(ext, 0) + 1

        weight = EXT_DEFAULTS.get(ext, 0.4)
        for cue, w in NAME_WEIGHTS:
            if cue in lower_name:
                weight = max(weight, w)

        weights_used[lower_name] = weight

        try:
            raw = f.read()
            text = ""
            if ext in TEXT_EXT:
                text = raw.decode("utf-8", errors="ignore")
            elif ext in DOCX_EXT:
                text = _extract_text_docx(raw)
            elif ext in PPTX_EXT:
                text = _extract_text_pptx(raw)
            elif ext in CSV_EXT:
                text = _extract_text_csv(raw, name)
            elif ext in PDF_EXT:
                # Extract PDF text and also store guidance hints for Value Managers
                text = _extract_text_pdf(raw)
                try:
                    hints = _pdf_review_hints(raw, name)
                except Exception:
                    hints = []
                if hints:
                    st.session_state["pdf_hints"][name] = hints
            else:
                text = f"[[FILE:{name}]]"

            if text.strip():
                chunks.append(f"\n# {name}\n{text.strip()}\n")
            else:
                chunks.append(f"\n# {name}\n[[NO-TEXT-EXTRACTED]]\n")
        except Exception:
            chunks.append(f"\n# {name}\n[[READ-ERROR]]\n")

    return "\n".join(chunks).strip(), counts, weights_used
    
# --------------- SME cues / analysis -----------------
FOUR_LEAF_KEYS: Dict[str, List[str]] = {
    "Human": [
        "team",
        "staff",
        "employee",
        "hire",
        "recruit",
        "training",
        "trained",
        "trainer",
        "onboarding",
        "mentor",
        "apprentice",
        "qualification",
        "certified",
        "cpd",
        "skills matrix",
        "safety training",
        "toolbox talk",
        "rota",
    ],
    "Structural": [
        # IAS 38 explicit / documented assets
        "contract",
        "agreement",
        "joint venture",
        "joint_venture",
        "jv",
        "mou",
        "grant",
        "licence",
        "license",
        "nda",
        "non-disclosure",
        "confidentiality",
        "ip register",
        "asset register",
        "register",
        "policy",
        "sop",
        "procedure",
        "process",
        "workflow",
        "protocol",
        "safety protocol",
        "process map",
        "manual",
        "guide",
        "handbook",
        "template",
        "checklist",
        "board pack",
        "board report",
        "pricing sheet",
        "tariff",
        "datasheet",
        "dataset",
        "qms",
        "iso 9001",
        "iso 27001",
        "knowledge base",
        "architecture",
        "crm",  # explicit system
    ],
    "Customer": [
        "client",
        "customer",
        "account",
        "lead",
        "opportunity",
        "pipeline",
        "quote",
        "proposal",
        "purchase order",
        "po",
        "invoice",
        "renewal",
        "retention",
        "distributor",
        "reseller",
        "channel",
        "customer success",
        "subscription",
    ],
    "Strategic Alliance": [
        "partner",
        "partnership",
        "alliance",
        "strategic",
        "mou",
        "joint venture",
        "framework agreement",
        "collaboration",
        "consortium",
        "university",
        "college",
        "council",
        "ngo",
        "integrator",
        "oem",
        "supplier agreement",
        "grant agreement",
        "licensor",
        "licensee",
        "jv",
    ],
}

TEN_STEPS = [
    "Identify",
    "Separate",
    "Protect",
    "Safeguard",
    "Manage",
    "Control",
    "Use",
    "Monitor",
    "Value",
    "Report",
]

SECTOR_CUES = {
    "GreenTech": [
        "recycling",
        "recycled",
        "waste",
        "biomass",
        "circular",
        "emissions",
        "co2e",
        "solar",
        "pv",
        "turbine",
        "kwh",
        "energy efficiency",
        "retrofit",
        "heat pump",
        "iso 14001",
        "esg",
        "sdg",
        "ofgem",
        "lca",
    ],
    "MedTech": [
        "iso 13485",
        "mhra",
        "ce mark",
        "clinical",
        "gcp",
        "patient",
        "medical device",
        "pms",
        "post-market surveillance",
    ],
    "AgriTech": [
        "soil",
        "irrigation",
        "seed",
        "fertiliser",
        "biomass",
        "yield",
        "farm",
    ],
}

# Qualitative market-growth and driver hints by sector.
# NOTE: Keys must match the sector labels used on page 1.

SECTOR_CUES = {
    ...
}

CAGR_API_BASE_URL = os.getenv("ICLICAI_CAGR_API_URL", "").strip()

def fetch_sector_cagr(sector: str) -> Optional[float]:
    ...
    # (block I gave you)
    ...

# Qualitative market-growth and driver hints by sector.
# NOTE: Keys must match the sector labels used on page 1.
SECTOR_CAGR_HINTS: Dict[str, str] = {
    ...
}

SECTOR_CAGR_HINTS: Dict[str, str] = {
    "GreenTech": (
        "GreenTech markets are expanding faster than the wider economy, driven by "
        "decarbonisation policies, energy-efficiency regulations, and demand for lower-cost, "
        "low-carbon solutions. Buyers are looking for technologies that reduce emissions, cut "
        "operating costs, and help them comply with ESG and climate-reporting requirements."
    ),
    "MedTech": (
        "MedTech demand is supported by ageing populations, pressure on healthcare budgets, "
        "and a shift toward outpatient and home-based care. Regulators are tightening quality "
        "and safety expectations, which favours solutions that demonstrate clinical evidence, "
        "regulatory compliance, and clear cost–benefit for hospitals and insurers."
    ),
    "AgriTech": (
        "AgriTech markets are growing as food-security concerns, climate change, and labour "
        "shortages push farmers toward data-driven and resource-efficient solutions. Adoption "
        "is strongest where tools can boost yields, reduce inputs like water and fertiliser, "
        "and provide traceability for export and premium markets."
    ),
    "Software / SaaS": (
        "Software and SaaS are benefiting from ongoing digitalisation, with organisations "
        "shifting from on-premise systems to cloud-based, subscription models. Buyers focus "
        "on time-to-value, ease of integration via APIs, cybersecurity, and predictable "
        "operating costs rather than large capital investments."
    ),
    "Professional & KIBS": (
        "Knowledge-intensive business services are growing as companies outsource specialised "
        "capabilities in analytics, compliance, digital transformation, and ESG. The market "
        "rewards firms that combine deep domain expertise with technology-enabled delivery and "
        "can demonstrate measurable impact on client performance and risk reduction."
    ),
    "Public / Gov / Non-profit": (
        "Public-sector and non-profit demand is shaped by long budget cycles, procurement "
        "rules, and policy priorities such as digital inclusion, resilience, and sustainability. "
        "Solutions that align with national strategies, offer strong value-for-money, and can be "
        "scaled across multiple agencies have a clear advantage."
    ),
    "_default": (
        "The company operates in a market that is expanding steadily as customers respond to "
        "regulatory change, digitalisation, and pressure to improve efficiency and resilience. "
        "Growth is strongest where providers can solve mission-critical problems, integrate with "
        "existing systems, and demonstrate clear economic and ESG benefits."
    ),
}

def get_sector_cagr_hint(sector: str) -> str:
    """
    Return a qualitative market-growth and driver narrative for the chosen sector.
    Falls back to '_default' if no exact match is found.
    """
    return SECTOR_CAGR_HINTS.get(sector, SECTOR_CAGR_HINTS["_default"])


def get_sector_market_context(sector: str) -> str:
    """
    Return a combined market-growth narrative for the chosen sector, using:
      - live CAGR from an external API (if available), and
      - qualitative drivers from SECTOR_CAGR_HINTS.
    """
    # Qualitative drivers
    narrative = SECTOR_CAGR_HINTS.get(sector, SECTOR_CAGR_HINTS["_default"])

    # Try to fetch a numeric CAGR from the external API
    cagr = fetch_sector_cagr(sector)
    if cagr is None:
        # No live data or API not configured: fall back to narrative only
        return narrative

    try:
        pct = round(float(cagr) * 100, 1)
    except Exception:
        # In case the API returns something odd
        return narrative

    return (
        f"Current external estimates suggest a compound annual growth rate of "
        f"around {pct}% per year for this sector. "
        f"{narrative}"
    )

def _strength_level(explicit_count: int, tacit_count: int) -> str:
    """
    Rough heuristic for capital strength based on artefact counts.
    """
    total = explicit_count + tacit_count
    if total >= 8:
        return "strong"
    if total >= 4:
        return "emerging"
    if total > 0:
        return "minimal"
    return "none"

def _step_band(score: int) -> str:
    """
    Map a Ten-Steps numeric score to a band for easier narrative.
    """
    if score >= 3:
        return "strong"
    if score == 2:
        return "developing"
    if score == 1:
        return "weak"
    return "gap"

def derive_vm_assumptions(
    sector: str,
    ic_summary: Dict[str, Dict[str, List[Any]]],
    ten_steps_scores: Dict[str, int],
) -> List[VMAssumption]:
    """
    Produce a list of narrative assumptions for the VM, derived from the analysis.
    The VM can accept or reject each assumption before it is written into the report.

    Expected ic_summary shape (per capital):
        ic_summary = {
            "Structural": {"explicit": [...], "tacit": [...]},
            "Human": {"explicit": [...], "tacit": [...]},
            "Customer": {"explicit": [...], "tacit": [...]},
            "Strategic": {"explicit": [...], "tacit": [...]},
        }

    Expected ten_steps_scores shape:
        ten_steps_scores = { "Identify": 2, "Separate": 1, ... }
    """

    assumptions: List[VMAssumption] = []

    # --- Capital strength diagnostics ---------------------------------------
    def _counts(cap: str) -> tuple[int, int]:
        cap_data = ic_summary.get(cap, {})
        return (
            len(cap_data.get("explicit", [])),
            len(cap_data.get("tacit", [])),
        )

    struct_exp, struct_tac = _counts("Structural")
    human_exp, human_tac = _counts("Human")
    cust_exp, cust_tac = _counts("Customer")
    strat_exp, strat_tac = _counts("Strategic")

    structural_level = _strength_level(struct_exp, struct_tac)
    human_level = _strength_level(human_exp, human_tac)
    customer_level = _strength_level(cust_exp, cust_tac)
    strategic_level = _strength_level(strat_exp, strat_tac)

    # --- Ten-Steps bands -----------------------------------------------------
    step_bands: Dict[str, str] = {
        step: _step_band(ten_steps_scores.get(step, 0))
        for step in TEN_STEPS
    }

    # --- 1. Market-growth / demand assumption -------------------------------
    market_context_text = get_sector_market_context(sector)
    assumptions.append(
        VMAssumption(
            key="market_pull",
            label="Market pull and growth context",
            narrative=(
                "We assume there is sustained market pull for this type of solution in the "
                f"chosen sector ({sector}). {market_context_text}"
            ),
            rationale=(
                "Based on the sector selected on page 1, external CAGR estimates (if "
                "available), and the qualitative SECTOR_CAGR_HINTS narrative."
            ),
            category="market",
            source_signals=[f"sector={sector}", "SECTOR_CAGR_HINTS", "CAGR_API"],
            confidence="medium",
        )
    )

    # --- 2. Innovation capacity assumption ----------------------------------
    innovative_capitals = sum(
        level in ("emerging", "strong")
        for level in (structural_level, human_level, strategic_level)
    )

    if innovative_capitals >= 2:
        innovation_narrative = (
            "The company appears to have a meaningful innovation base, with core know-how, "
            "people, and partnering assets that can support repeatable development of new "
            "solutions and service improvements."
        )
        innovation_conf = "high"
    else:
        innovation_narrative = (
            "The company seems to be at an earlier stage in building its innovation base, "
            "with some promising assets but still limited depth in documented processes, "
            "skills, or strategic partnerships."
        )
        innovation_conf = "medium"

    assumptions.append(
        VMAssumption(
            key="innovation_capacity",
            label="Innovation capacity",
            narrative=innovation_narrative,
            rationale=(
                "Derived from the relative strength of Structural, Human, and Strategic "
                "capitals (tacit + explicit artefact counts)."
            ),
            category="innovation",
            source_signals=[
                f"Structural={structural_level}",
                f"Human={human_level}",
                f"Strategic={strategic_level}",
            ],
            confidence=innovation_conf,
        )
    )

    # --- 3. Commercial readiness assumption ---------------------------------
    use_band = step_bands.get("Use", "gap")
    monitor_band = step_bands.get("Monitor", "gap")
    value_band = step_bands.get("Value", "gap")

    if customer_level in ("emerging", "strong") and use_band in ("developing", "strong"):
        comm_narrative = (
            "We assume the company is beyond initial proof-of-concept and is already "
            "testing or deploying its solution with real customers, with a pathway toward "
            "repeatable commercial delivery."
        )
        comm_conf = "high" if value_band in ("developing", "strong") else "medium"
    else:
        comm_narrative = (
            "We assume the company is still in earlier commercial-validation stages and "
            "needs to strengthen its customer base and repeatable commercial model before "
            "scaling."
        )
        comm_conf = "medium"

    assumptions.append(
        VMAssumption(
            key="commercial_readiness",
            label="Commercial readiness",
            narrative=comm_narrative,
            rationale=(
                "Based on Customer Capital strength and the Use / Monitor / Value steps in "
                "the Ten-Steps analysis."
            ),
            category="market",
            source_signals=[
                f"Customer={customer_level}",
                f"Use_step={use_band}",
                f"Monitor_step={monitor_band}",
                f"Value_step={value_band}",
            ],
            confidence=comm_conf,
        )
    )

    # --- 4. IP & governance maturity assumption -----------------------------
    identify_band = step_bands.get("Identify", "gap")
    protect_band = step_bands.get("Protect", "gap")
    safeguard_band = step_bands.get("Safeguard", "gap")
    manage_band = step_bands.get("Manage", "gap")
    control_band = step_bands.get("Control", "gap")

    ip_bands = [identify_band, protect_band, safeguard_band, manage_band, control_band]
    strong_ip = sum(b in ("developing", "strong") for b in ip_bands)

    if strong_ip >= 3:
        ip_narrative = (
            "We assume the company has at least basic IP and knowledge-governance processes "
            "in place, with identified core assets and some level of protection, "
            "management, and access control."
        )
        ip_conf = "high"
    else:
        ip_narrative = (
            "We assume core IP and knowledge-governance processes are still emerging, with "
            "gaps in how assets are identified, protected, and controlled across the "
            "organisation."
        )
        ip_conf = "medium"

    assumptions.append(
        VMAssumption(
            key="ip_governance_maturity",
            label="IP and governance maturity",
            narrative=ip_narrative,
            rationale=(
                "Derived from the Identify / Protect / Safeguard / Manage / Control steps "
                "in the Ten-Steps analysis."
            ),
            category="ten-steps",
            source_signals=[
                f"Identify={identify_band}",
                f"Protect={protect_band}",
                f"Safeguard={safeguard_band}",
                f"Manage={manage_band}",
                f"Control={control_band}",
            ],
            confidence=ip_conf,
        )
    )

    # --- 5. Execution risk assumption ---------------------------------------
    weak_or_gap = sum(b in ("weak", "gap") for b in step_bands.values())
    if weak_or_gap >= 5:
        risk_narrative = (
            "We assume there is a material execution risk: several foundational activities "
            "in the asset lifecycle are weak or missing, which could slow delivery, weaken "
            "negotiating power, or block investment until addressed."
        )
        risk_conf = "medium"
    else:
        risk_narrative = (
            "We assume execution risk is manageable: there are still gaps, but the company "
            "has enough structure in place to support growth if the most critical steps are "
            "prioritised in the next 12–24 months."
        )
        risk_conf = "medium"

    assumptions.append(
        VMAssumption(
            key="execution_risk",
            label="Execution risk profile",
            narrative=risk_narrative,
            rationale=(
                "Based on the distribution of strong vs. weak/gap scores across all Ten Steps."
            ),
            category="ten-steps",
            source_signals=[f"weak_or_gap_steps={weak_or_gap}"],
            confidence=risk_conf,
        )
    )

    return assumptions
    
# Explicit structural cues (IAS 38-compliant artefact hints)
EXPLICIT_STRUCTURAL_CUES: List[str] = [
    "contract",
    "agreement",
    "msa",
    "sow",
    "sla",
    "mou",
    "joint venture",
    "joint_venture",
    "jv",
    "grant",
    "licence",
    "license",
    "ip register",
    "asset register",
    "register",
    "policy",
    "sop",
    "procedure",
    "process",
    "workflow",
    "protocol",
    "safety protocol",
    "process map",
    "manual",
    "guide",
    "handbook",
    "template",
    "checklist",
    "board pack",
    "board report",
    "pricing sheet",
    "tariff",
    "datasheet",
    "dataset",
    "qms",
    "iso 9001",
    "iso 27001",
    "crm",
]

# ESG & Seven Stakeholder cues (Sugai / Weir)
ESG_CUES: List[str] = [
    "esg",
    "sdg",
    "carbon",
    "emissions",
    "net zero",
    "scope 1",
    "scope 2",
    "scope 3",
    "sustainability",
    "governance",
    "diversity",
    "inclusion",
    "impact",
    "stakeholder",
]

SEVEN_STAKEHOLDER_CUES: List[str] = [
    "employee",
    "staff",
    "worker",
    "investor",
    "shareholder",
    "lender",
    "customer",
    "client",
    "supplier",
    "vendor",
    "partner",
    "alliance",
    "community",
    "local authority",
    "municipality",
    "ngo",
    "nature",
    "environment",
    "biodiversity",
]

# --------------- ANALYSIS ENGINE ---------------------
def _analyse_weighted(
    text: str,
    weights_by_file: Dict[str, float],
) -> Tuple[Dict[str, Any], Dict[str, float], Dict[str, Any], int]:
    """
    Weighted Four-Leaf & Ten-Steps.
    Returns:
      ic_map (with tick/narrative/score),
      leaf_scores (raw weighted scores for 4-leaf),
      ten (scores+narratives),
      quality% (heuristic)
    """
    sector = st.session_state.get("sector", "Other")
    t_all = (text or "").lower()

    leaf_scores: Dict[str, float] = {
        "Human": 0.0,
        "Structural": 0.0,
        "Customer": 0.0,
        "Strategic Alliance": 0.0,
    }
    step_scores: Dict[str, float] = {s: 0.0 for s in TEN_STEPS}

    sector_present = False
    if sector in SECTOR_CUES:
        if any(c in t_all for c in SECTOR_CUES[sector]):
            sector_present = True

    # ----- Structural vs Tacit weighting -----
    max_weight = max(weights_by_file.values() or [0.4])

    # Base structural emphasis from explicit cues anywhere in the text (IAS 38 explicit assets)
    for cue in EXPLICIT_STRUCTURAL_CUES:
        if cue in t_all:
            leaf_scores["Structural"] += max_weight * 1.5  # audit-ready bump

    # Four-Leaf cues (with sector reinforcement)
    for leaf, cues in FOUR_LEAF_KEYS.items():
        eff = list(cues)
        if sector_present and leaf in ("Structural", "Customer", "Strategic Alliance"):
            eff += SECTOR_CUES[sector]
        base = 0.0
        for cue in eff:
            if cue in t_all:
                base += max_weight
        leaf_scores[leaf] += base

    # Ten-Steps scoring (file-name based + ESG / FRAND cues)
    def bump(step: str, amt: float) -> None:
        step_scores[step] = step_scores.get(step, 0.0) + amt

    for fname, w in (weights_by_file or {}).items():
        n = fname.lower()

        # Contracts / grants / agreements → Structural dominates, Customer/SA secondary
        if any(k in n for k in ["contract", "msa", "sow", "sla", "po", "agreement"]):
            leaf_scores["Structural"] += 2.5 * w
            leaf_scores["Customer"] += 1.0 * w
            bump("Control", 2.0 * w)
            bump("Use", 2.5 * w)

        if any(k in n for k in ["joint_venture", "joint venture", "jv", "mou", "grant"]):
            leaf_scores["Structural"] += 2.5 * w
            leaf_scores["Strategic Alliance"] += 1.5 * w
            bump("Control", 2.0 * w)
            bump("Use", 2.0 * w)

        # Knowledge / SOP / KMP / safety / ISO → Structural + Human
        if any(k in n for k in ["knowledge", "kmp", "sop", "process", "safety", "protocol", "risk", "qms", "iso"]):
            leaf_scores["Structural"] += 1.8 * w
            leaf_scores["Human"] += 0.8 * w
            bump("Identify", 1.8 * w)
            bump("Separate", 1.4 * w)
            bump("Manage", 1.6 * w)
            bump("Safeguard", 1.0 * w)

        # Specs/slides/canvas → Structural + Use
        if any(k in n for k in ["spec", "canvas", "deck", "slides", "pptx"]):
            leaf_scores["Structural"] += 0.8 * w
            bump("Identify", 0.8 * w)
            bump("Use", 0.6 * w)

        # Pricing/licensing hints → Use/Value (multi value streams)
        if any(k in n for k in ["price", "pricing", "royalty", "subscription", "oem", "white label"]):
            bump("Use", 1.2 * w)
            bump("Value", 1.6 * w)

        # Governance/reporting → Monitor/Report
        if any(k in n for k in ["board", "report", "dashboard", "audit"]):
            bump("Report", 1.4 * w)
            bump("Monitor", 1.2 * w)

    # ESG & Seven Stakeholder presence → boost Report/Value (double materiality)
    esg_hits = any(c in t_all for c in ESG_CUES)
    stakeholder_hits = any(c in t_all for c in SEVEN_STAKEHOLDER_CUES)
    if esg_hits or stakeholder_hits:
        bump("Report", 1.2)
        bump("Value", 1.0)

    if sector_present:
        bump("Use", 0.8)
        bump("Report", 0.5)

    # Make sure Structural "wins" when explicit + tacit both present:
    # if Structural>0 and (Customer or SA also high), add a small dominance bump.
    if leaf_scores["Structural"] > 0 and (leaf_scores["Customer"] > 0 or leaf_scores["Strategic Alliance"] > 0):
        leaf_scores["Structural"] *= 1.15  # dominance tweak

    # Convert leaf_scores -> ticks & narratives
    ic_map: Dict[str, Any] = {}
    avg_leaf = (sum(leaf_scores.values()) / max(1, len(leaf_scores)))
    threshold = max(1.0, avg_leaf * 0.6)

    for leaf, score in leaf_scores.items():
        tick = score >= threshold
        if leaf == "Human":
            nar = (
                "Human Capital evidenced (values, awards, training and safety practice), but competency and role-mapping "
                "should be consolidated into a formal skills register."
                if tick
                else "Human Capital is not yet clearly evidenced; competency mapping, training logs and safety records are needed."
            )
        elif leaf == "Structural":
            nar = (
                "Structural Capital appears IAS 38-ready in places (contracts, SOPs, protocols, registers, CRM and board packs "
                "are present), supporting audit-ready recognition on the balance sheet."
                if tick
                else "Structural Capital is under-documented; explicit artefacts (contracts, registers, SOPs, board packs, "
                "pricing, datasets, CRM) should be consolidated into an auditable IA Register."
            )
        elif leaf == "Customer":
            nar = (
                "Customer Capital is evidenced through relationships, renewal logic and channels, supporting recurring value capture "
                "and future licensing opportunities."
                if tick
                else "Customer Capital appears weak in the evidence; relationship histories, renewals, CRM and pipeline data should be structured."
            )
        else:
            nar = (
                "Strategic Alliance Capital is evidenced (JVs, MoUs, partners, universities, councils), enabling co-creation "
                "and ecosystem-based licensing opportunities."
                if tick
                else "Strategic alliances are not clearly evidenced; JV/MoU documentation and partner frameworks are needed."
            )
        ic_map[leaf] = {"tick": tick, "narrative": nar, "score": round(score, 2)}

    # Ten-Steps scores
    base = 3.0
    ten_scores: List[int] = []
    ten_narrs: List[str] = []

    for step in TEN_STEPS:
        s_float = base + step_scores.get(step, 0.0)
        s = int(max(1, min(10, round(s_float))))
        ten_scores.append(s)
        ten_narrs.append(f"{step}: readiness ≈ {s}/10.")

    ten = {"scores": ten_scores, "narratives": ten_narrs}

    # Evidence quality metric
    files_factor = min(1.0, len(weights_by_file) / 6.0)
    leaf_div = sum(1 for v in ic_map.values() if v["tick"]) / 4.0
    weight_mean = (sum(weights_by_file.values()) / max(1, len(weights_by_file))) if weights_by_file else 0.4
    quality = int(round(100 * (0.45 * files_factor + 0.35 * leaf_div + 0.20 * min(1.0, weight_mean))))

    return ic_map, leaf_scores, ten, quality

# --------------- INTERPRETIVE NARRATIVE --------------
def _build_interpreted_summary(
    case: str,
    leaf_scores: Dict[str, float],
    ic_map: Dict[str, Any],
    ten: Dict[str, Any],
    evidence_quality: int,
    context: Dict[str, str],
) -> str:
    sector = st.session_state.get("sector", "Other")
    size = st.session_state.get("company_size", "Micro (1–10)")

    strengths = [k for k, v in ic_map.items() if v.get("tick")]
    gaps = [k for k, v in ic_map.items() if not v.get("tick")]

    ts = ten.get("scores") or [5] * len(TEN_STEPS)
    strong_steps = [s for s, sc in zip(TEN_STEPS, ts) if sc >= 7]
    weak_steps = [s for s, sc in zip(TEN_STEPS, ts) if sc <= 5]

    # Detect whether ESG & Seven Stakeholder cues are present
    narrative_text = context.get("why", "") + " " + context.get("markets", "")
    seven_hit = any(c in narrative_text.lower() for c in SEVEN_STAKEHOLDER_CUES)
    esg_hit = any(c in narrative_text.lower() for c in ESG_CUES)

    # Quick handle for Structural strength and evidence depth
    structural_row = ic_map.get("Structural", {})
    structural_tick = bool(structural_row.get("tick"))
    evidence_thin = evidence_quality < 50  # tweak threshold later if needed

    # 1) Context & positioning
    p1 = (
        f"{case} is a {size} in {sector}. Based on uploaded artefacts and company context, the company shows an "
        f"emerging ability to codify and scale its operating model, with measurable signals across "
        f"{', '.join(strengths) if strengths else 'selected IC dimensions'}."
    )

    # 2) Four-Leaf interpretation (explicit vs tacit, IAS 38)
    if strengths:
        p2a = f"Strengths concentrate in {', '.join(strengths)}" + (f"; gaps are {', '.join(gaps)}." if gaps else ".")
    else:
        p2a = "Strengths are not yet well-evidenced; additional artefacts are required."

    if structural_tick and not evidence_thin:
        # Strong wording only when Structural is ticked AND evidence depth is reasonable
        p2b = (
            "Evidence points to maturing Structural Capital where explicit artefacts — contracts, SOPs, protocols, registers, "
            "board materials, CRM and datasets — are present. These are the primary candidates for IAS 38-compliant recognition on "
            "the balance sheet. Human, Customer and Strategic Alliance Capital are reflected mainly through tacit know-how, "
            "relationships and informal practice, which require codification before they become audit-ready."
        )
    else:
        # Softer wording when evidence is thin or Structural is not clearly strong
        p2b = (
            "At this stage Structural Capital is only weakly evidenced. There are early or fragmentary signals (for example references "
            "to contracts, processes, governance or data), but not yet in a form that could support IAS 38 recognition. The immediate "
            "priority is to consolidate existing documents into a simple IA Register and to separate clearly what is codified from what "
            "still sits in people’s heads and informal practice across Human, Customer and Strategic Alliance Capital."
        )

    p2 = p2a + " " + p2b

    # 3) Ten-Steps insight and readiness for licensing
    if strong_steps or weak_steps:
        p3 = (
            f"Ten-Steps patterns indicate strong {', '.join(strong_steps) if strong_steps else 'foundations'}; "
            f"progress is constrained by {', '.join(weak_steps) if weak_steps else 'later-stage governance, valuation and reporting readiness'}."
        )
    else:
        p3 = (
            "Ten-Steps scores suggest a developing baseline; company-side review will refine scoring as artefacts are "
            "consolidated and IA governance is embedded."
        )

    # 4) Seven Stakeholder / ESG framing and value streams
    if seven_hit or esg_hit:
        p4_intro = (
            "Using the Seven Stakeholder Model (SSM) as defined by Professor Philip Sugai and Dr Maria Weir, "
            "the company can frame value creation across employees, investors, customers, partners and suppliers, "
            "communities and the natural environment. This provides a structured way to connect ESG performance "
            "and double materiality to concrete intangible assets."
        )
    else:
        p4_intro = (
            "Applying the Seven Stakeholder Model (SSM) as defined by Professor Philip Sugai and Dr Maria Weir "
            "would allow the company to frame value creation across employees, investors, customers, partners and "
            "suppliers, communities and the natural environment, even if these links are not yet fully articulated "
            "in the evidence set."
        )

    p4_mid = (
        "From a commercialisation perspective, explicit Structural Capital (contracts, data, software, methods, indices, "
        "protocols) can support multiple simultaneous value streams — including revenue licences, access or community "
        "licences, co-creation arrangements and data/algorithm sharing — provided that ownership, rights and governance "
        "are clarified."
    )

    actions = [
        "Create a single IA Register linking all explicit artefacts (contracts, JVs, SOPs, protocols, datasets, board packs, CRM).",
        "Map each explicit asset to at least one licensing-ready value stream (revenue, access/community, co-creation, defensive or data/algorithm sharing).",
        "Introduce quarterly governance reporting (board pack + KPI dashboard) to strengthen Monitor and Report and to evidence ESG and stakeholder impacts.",
        "Define valuation approach (IAS 38 fair value) and link to licensing templates so that audit-ready Structural Capital supports near-term monetisation.",
        "Formalise competency matrices and training logs so that tacit Human Capital can be progressively codified into Structural Capital.",
    ]
    p4_actions = "Assumptions & Action Plan:\n" + "\n".join([f"• {a}" for a in actions])

    p4 = p4_intro + "\n\n" + p4_mid + "\n\n" + p4_actions

    # 5) Evidence quality and next evidence requests
    missing = (
        "Request additional artefacts: CRM/renewal data, NDA/licence/royalty and access terms, IA or IP registers, "
        "board/management reports, and any ESG or stakeholder dashboards used in internal decision-making."
    )
    p5 = f"Evidence quality ≈ {evidence_quality}% (heuristic). {missing}"

    return "\n\n".join([p1, p2, p3, p4, p5])

# --------- COMPANY CONTEXT AUTO-SPLIT HELPER ----------
def _auto_split_expert_block(text: str) -> Dict[str, str]:
    """
    Take a single pasted block and try to split it across:
    why_service, stage, plan_s, plan_m, plan_l, markets_why, sale_price_why
    using blank lines or sentence boundaries.
    """
    t = (text or "").strip()
    keys = [
        "why_service",
        "stage",
        "plan_s",
        "plan_m",
        "plan_l",
        "markets_why",
        "sale_price_why",
    ]
    if not t:
        return {k: "" for k in keys}

    blocks = [b.strip() for b in t.replace("\r\n", "\n").split("\n\n") if b.strip()]

    if len(blocks) < 5:
        blocks = [s.strip() for s in re.split(r"(?<=[.!?])\s+", t) if s.strip()]

    out: Dict[str, str] = {k: "" for k in keys}
    for k, chunk in zip(keys, blocks):
        out[k] = chunk
    return out

# ------------------ SESSION DEFAULTS -----------------
ss = st.session_state
ss.setdefault("case_name", "Untitled Company")
ss.setdefault("company_size", "Micro (1–10)")
ss.setdefault("sector", "Other")
ss.setdefault("uploads", [])
ss.setdefault("combined_text", "")
ss.setdefault("ic_map", {})
ss.setdefault("ten_steps", {})
ss.setdefault("narrative", "")
ss.setdefault("leaf_scores", {})
ss.setdefault("evidence_quality", 0)
ss.setdefault("file_counts", {})

# Expert prompts
ss.setdefault("why_service", "")
ss.setdefault("stage", "")
ss.setdefault("plan_s", "")
ss.setdefault("plan_m", "")
ss.setdefault("plan_l", "")
ss.setdefault("markets_why", "")
ss.setdefault("sale_price_why", "")
ss.setdefault("full_context_block", "")
ss.setdefault("auto_split_on_save", True)

# LIP Assistant state
ss.setdefault("lip_history", [])

# Asset verification (overall notes)
ss.setdefault("verification_notes", "")

SIZES = [
    "Micro (1–10)",
    "Small (11–50)",
    "Medium (51–250)",
    "Large (250+)",
]
SECTORS = [
    "Food & Beverage",
    "MedTech",
    "GreenTech",
    "AgriTech",
    "Biotech",
    "Software/SaaS",
    "FinTech",
    "EdTech",
    "Manufacturing",
    "Creative/Digital",
    "Professional Services",
    "Mobility/Transport",
    "Energy",
    "Other",
]
# --------------- SIDEBAR BRANDING & NAV ---------------
with st.sidebar:
    # IMPAC3T-IP logo (top) – safe load, moderate width
    try:
        if IMPACT3T_LOGO_PATH.is_file():
            st.image(str(IMPACT3T_LOGO_PATH), width=170)
        else:
            st.markdown("**IMPAC3T-IP**")
    except Exception:
        st.markdown("**IMPAC3T-IP**")

    st.markdown("---")

    # Navigation (middle of sidebar)
    st.markdown("### Navigate")
    page = st.radio(
        "Go to:",
        [
            "Company",
            "Analyse Evidence",
            "Asset Verification",
            "LIP Console",
            "Reports",
            "Licensing Templates",
            "LIP Assistant",
        ],
        index=0,
    )

    st.markdown("---")

    # EU flag + funding line – **footer**
    try:
        if EU_FLAG_PATH.is_file():
            st.image(str(EU_FLAG_PATH), width=80)
        else:
            st.markdown("EU-funded tool")
    except Exception:
        st.markdown("EU-funded tool")

    st.caption(
        "This tool has been developed within the IMPAC3T-IP project, which has received "
        "funding from the European Union's Horizon Europe programme under Grant Agreement "
        "No. 101135832."
    )
    
# -------------------- PAGES -------------------------
  
# 1) COMPANY (with required prompts + auto-split)
if page == "Company":
    st.header("Company details")

    # ---------------- COMPANY FORM ----------------
    with st.form("company_form"):
        c1, c2, c3 = st.columns([1.1, 1, 1])
        with c1:
            case_name = st.text_input("Company name *", ss.get("case_name", ""))
        with c2:
            size = st.selectbox(
                "Company size",
                SIZES,
                index=SIZES.index(ss.get("company_size", SIZES[0])),
            )
        with c3:
            current_sector = ss.get("sector", "Other")
            sector_index = (
                SECTORS.index(current_sector)
                if current_sector in SECTORS
                else SECTORS.index("Other")
            )
            sector = st.selectbox("Sector / Industry", SECTORS, index=sector_index)

        st.markdown("#### Company context (required)")
        full_block = st.text_area(
            "Optional: paste full context here (one block, then auto-fill below)",
            ss.get("full_context_block", ""),
            help=(
                "You can paste your whole narrative here once, then tick auto-split and the tool will "
                "try to populate the individual questions for you."
            ),
            height=80,
        )
        auto_split = st.checkbox(
            "Auto-split pasted context into fields on Save",
            value=ss.get("auto_split_on_save", True),
            help=(
                "If ticked, the block above will be split across the questions below "
                "when you click Save."
            ),
        )

        why_service = st.text_area(
            "Why is the company seeking this service? *",
            ss.get("why_service", ""),
            height=90,
        )
        stage = st.text_area(
            "What stage are the products/services at? *",
            ss.get("stage", ""),
            height=90,
        )

        c4, c5, c6 = st.columns(3)
        with c4:
            plan_s = st.text_area(
                "3a) Short-term plan (0–6m) *",
                ss.get("plan_s", ""),
                height=90,
            )
        with c5:
            plan_m = st.text_area(
                "3b) Medium-term plan (6–24m) *",
                ss.get("plan_m", ""),
                height=90,
            )
        with c6:
            plan_l = st.text_area(
                "3c) Long-term plan (24m+) *",
                ss.get("plan_l", ""),
                height=90,
            )

        markets_why = st.text_area(
            "4) Which markets fit and why? *",
            ss.get("markets_why", ""),
            height=90,
        )
        sale_price_why = st.text_area(
            "5) If selling tomorrow, target price & why? *",
            ss.get("sale_price_why", ""),
            height=90,
        )

        st.caption(
            "Uploads are held in session until analysis. Nothing is written to server until export."
        )
        uploads = st.file_uploader(
            "Upload evidence (PDF, DOCX, TXT, CSV, XLSX, PPTX, images)",
            type=[
                "pdf",
                "docx",
                "txt",
                "csv",
                "xlsx",
                "pptx",
                "png",
                "jpg",
                "jpeg",
                "webp",
            ],
            accept_multiple_files=True,
            key="uploader_main",
        )

        submitted = st.form_submit_button("Save details")

        if submitted:
            # Optional auto-split of a single pasted narrative into all questions
            if auto_split and full_block.strip():
                derived = _auto_split_expert_block(full_block)
                if not why_service.strip() and derived["why_service"]:
                    why_service = derived["why_service"]
                if not stage.strip() and derived["stage"]:
                    stage = derived["stage"]
                if not plan_s.strip() and derived["plan_s"]:
                    plan_s = derived["plan_s"]
                if not plan_m.strip() and derived["plan_m"]:
                    plan_m = derived["plan_m"]
                if not plan_l.strip() and derived["plan_l"]:
                    plan_l = derived["plan_l"]
                if not markets_why.strip() and derived["markets_why"]:
                    markets_why = derived["markets_why"]
                if not sale_price_why.strip() and derived["sale_price_why"]:
                    sale_price_why = derived["sale_price_why"]

            missing = [
                ("Company name", case_name),
                ("Why service", why_service),
                ("Stage", stage),
                ("Short plan", plan_s),
                ("Medium plan", plan_m),
                ("Long plan", plan_l),
                ("Markets & why", markets_why),
                ("Sale price & why", sale_price_why),
            ]
            missing_fields = [
                label for (label, val) in missing if not (val or "").strip()
            ]
            if missing_fields:
                st.error(
                    "Please complete required fields: " + ", ".join(missing_fields)
                )
            else:
                ss["case_name"] = case_name
                ss["company_size"] = size
                ss["sector"] = sector
                ss["why_service"] = why_service.strip()
                ss["stage"] = stage.strip()
                ss["plan_s"] = plan_s.strip()
                ss["plan_m"] = plan_m.strip()
                ss["plan_l"] = plan_l.strip()
                ss["markets_why"] = markets_why.strip()
                ss["sale_price_why"] = sale_price_why.strip()
                ss["full_context_block"] = full_block
                ss["auto_split_on_save"] = auto_split
                if uploads:
                    ss["uploads"] = uploads
                st.success("Saved company details & context.")

    # ---------------- POST-FORM INFO ----------------
    if ss.get("uploads"):
        st.info(
            f"{len(ss['uploads'])} file(s) stored in session. "
            "Go to **Analyse Evidence** next."
        )

    # ---------------- VM TOOLBOX ----------------
    st.markdown("---")
    st.subheader("Value Manager toolbox (evidence preparation)")

    with st.expander("1. Convert PDFs and spreadsheets before upload"):
        st.markdown(
            """
Use these external tools **before** uploading evidence, especially for large reports:

- **PDF → Word (DOCX)**: Convert long ESG or strategy PDFs into Word so key text can be pasted into the Company context fields and uploaded as DOCX.
- **Excel → CSV**: Export key sheets (contracts, IP registers, ESG KPIs) as CSV so IC-LicAI can read headers and sample rows.

Remember: never upload anything that breaches confidentiality or data-sharing rules agreed with the client.
"""
        )
        st.markdown(
            """
**Suggested actions for VMs:**

1. If a PDF is mostly tables/diagrams with little text, create a short DOCX summary capturing:
   - main KPIs,
   - key contracts / partners,
   - any explicit IP / IA registers,
   - ESG claims that matter for value.
2. Save that summary as a DOCX and upload it alongside the original PDF.
3. For complex Excel files, export the most relevant sheet as CSV and upload that CSV.
"""
        )

    with st.expander("2. Optional AI helper for drafting the Company context"):
        st.markdown(
            """
You can use an AI assistant (e.g. ChatGPT) **before** opening this tool to help draft the answers for:

- Why is the company seeking this service?
- Stage of products/services
- Short / medium / long-term plans
- Markets & why
- Target sale price & why

**Suggested prompt to paste into your AI tool:**
"""
        )
        st.code(
            "You are helping me prepare an intellectual capital and licensing diagnostic for a company. "
            "Based on the attached business plan or PDF, draft:\n"
            "1) Why the company is seeking this service (funding, growth, governance, ESG, etc.),\n"
            "2) The stage of each key product/service,\n"
            "3) Short (0–6m), medium (6–24m) and long-term (24m+) plans,\n"
            "4) Which markets fit and why (sectors, partners, channels, not just countries),\n"
            "5) If the owners wanted to sell tomorrow, a realistic target sale price and the reasons.\n\n"
            "Write each as a clear paragraph I can paste into a diagnostic tool. Do not invent numbers; "
            "only use figures and evidence you can see in the document."
        )
        st.markdown(
            """
After AI has drafted this, **you remain responsible** for checking accuracy against the original documents
and adjusting anything that does not match reality.
"""
        )

    with st.expander("3. IC System View note (complex systems lens)"):
        st.markdown(
            """
Before running analysis, create a 1–2 page note (IC_SystemView_<CompanyName>.txt) with:

1. **Market system & actors** – who is upstream (suppliers, data sources, funders) and downstream (customers, end-users, regulators, communities)?
2. **Innovation loops & learning** – how products/services are improved; pilots, feedback, R&D, training.
3. **Key intangible / data / IP assets** – contracts, protocols, datasets, algorithms, indices, software, brand, training packages.
4. **Current & potential value streams** – who pays, who accesses for free or reduced price, who benefits indirectly.
5. **Key tensions & risks** – ESG tensions, dependency on single partners, regulatory or reputational risks.

Upload this note **as evidence** together with the other files. It gives IC-LicAI and the ILP a richer starting point.
"""
        )
        
# 2) ANALYSE EVIDENCE (with radar / evidence quality)
elif page == "Analyse Evidence":
    st.header("Evidence Dashboard & Analysis")

    col1, col2 = st.columns([1, 1])

    # ------------ LEFT COLUMN: EVIDENCE QUALITY ------------
    with col1:
        st.subheader("Evidence Quality")
        eq = int(ss.get("evidence_quality", 0))
        st.progress(min(100, max(0, eq)) / 100.0)
        st.caption(
            f"{eq}% coverage (heuristic — based on artefact mix and IC diversity)."
        )

        counts = ss.get("file_counts", {}) or {}
        if counts:
            st.markdown("**Files by type (session):**")
            for ext, n in counts.items():
                st.markdown(f"- `{ext}` → {n} file(s)")
        else:
            st.caption("No files analysed yet.")

    # ------------ RIGHT COLUMN: RADAR + TEN-STEPS ------------
    with col2:
        st.subheader("IC Radar (4-Leaf + Ten-Steps)")

        ic_map: Dict[str, Any] = ss.get("ic_map", {})
        ten = ss.get(
            "ten_steps",
            {"scores": [5] * len(TEN_STEPS), "narratives": [f"{s}: tbd" for s in TEN_STEPS]},
        )

        leaf_labels = ["Human", "Structural", "Customer", "Strategic Alliance"]
        leaf_vals = [float(ic_map.get(l, {}).get("score", 0.0)) for l in leaf_labels]

        # Gate the radar by evidence quality so we don't over-interpret tiny evidence sets
        eq_local = int(ss.get("evidence_quality", 0))

        if eq_local < 35 or not any(v > 0 for v in leaf_vals):
            st.caption(
                "Radar will appear once evidence coverage is above ~35% and IC signals are detected. "
                "At low coverage the map would be misleading, so treat current results as a very early scan."
            )
        else:
            fig_leaf = go.Figure()
            fig_leaf.add_trace(
                go.Scatterpolar(
                    r=leaf_vals + leaf_vals[:1],
                    theta=leaf_labels + leaf_labels[:1],
                    fill="toself",
                    name="IC Intensity",
                )
            )
            fig_leaf.update_layout(
                polar=dict(
                    radialaxis=dict(
                        visible=True,
                        range=[0, max(leaf_vals) or 1],
                    )
                ),
                showlegend=False,
                margin=dict(l=20, r=20, t=20, b=20),
            )
            st.plotly_chart(fig_leaf, use_container_width=True)

        step_scores = ten.get("scores") or [5] * len(TEN_STEPS)
        if step_scores:
            fig_steps = go.Figure(
                data=[
                    go.Bar(
                        x=TEN_STEPS,
                        y=step_scores,
                    )
                ]
            )
            fig_steps.update_layout(
                yaxis=dict(range=[0, 10]),
                margin=dict(l=20, r=20, t=20, b=40),
            )
            st.plotly_chart(fig_steps, use_container_width=True)

    # ------------ INTERPRETED SUMMARY TEXT AREA ------------
    st.markdown("---")
    st.subheader("Interpreted Analysis Summary (first 5000 chars)")
    combined = ss.get("combined_text", "")
    st.text_area(
        "Summary view (read-only; editable in LIP Console)",
        combined[:5000],
        height=260,
        key="combined_preview",
    )

    # ------------ RUN ANALYSIS BUTTON ------------
    if st.button("Run analysis now"):
        uploads: List[Any] = ss.get("uploads") or []
        extracted, counts, weights = _read_text(uploads)

        ss["file_counts"] = counts or {}

        context = {
            "why": ss.get("why_service", ""),
            "stage": ss.get("stage", ""),
            "plan_s": ss.get("plan_s", ""),
            "plan_m": ss.get("plan_m", ""),
            "plan_l": ss.get("plan_l", ""),
            "markets": ss.get("markets_why", ""),
            "sale": ss.get("sale_price_why", ""),
        }

        context_stub = (
            f"[CTX] why={context['why'][:140]} | stage={context['stage'][:140]} | "
            f"plans=({context['plan_s'][:60]}/{context['plan_m'][:60]}/{context['plan_l'][:60]}) | "
            f"markets={context['markets'][:140]} | sale={context['sale'][:60]}"
        )

        combined_text_for_detection = (extracted + "\n\n" + context_stub).strip().lower()

        ic_map, leaf_scores, ten, quality = _analyse_weighted(
            combined_text_for_detection,
            weights,
        )

        case = ss.get("case_name", "Untitled Company")
        interpreted = _build_interpreted_summary(
            case,
            leaf_scores,
            ic_map,
            ten,
            quality,
            context,
        )

        ss["combined_text"] = interpreted
        ss["ic_map"] = ic_map
        ss["ten_steps"] = ten
        ss["leaf_scores"] = leaf_scores
        ss["evidence_quality"] = quality

        if len(extracted.strip()) < 100:
            st.warning(
                "Little machine-readable text was extracted (DOCX/PPTX/CSV/CSV extraction is enabled). "
                "If PDFs dominate, consider adding a brief TXT note or exporting key pages to DOCX."
            )

        st.success("Analysis complete. Open **LIP Console** to review the summary and IC map.")

    # ------------ PDF REVIEW HINTS FOR VALUE MANAGERS ------------
    uploads = ss.get("uploads", [])
    pdf_files = [f for f in uploads or [] if getattr(f, "name", "").lower().endswith(".pdf")]

    if pdf_files:
        st.markdown("---")
        st.subheader("Suggested PDF pages to review (beta)")

        st.caption(
            "These hints scan each PDF for pages that mention tables, KPIs, IP, contracts, "
            "markets or technology. They are **for human review only** – they do not change "
            "scores or assumptions. Use them to jump to the most IC-relevant pages in the PDF."
        )

        for f in pdf_files:
            fname = getattr(f, "name", "unnamed.pdf")
            try:
                raw = f.getvalue() if hasattr(f, "getvalue") else f.read()
            except Exception:
                raw = None

            if not raw:
                continue

            hints = _pdf_review_hints(raw, fname)

            st.markdown(f"**{fname}**")
            if hints:
                for h in hints:
                    st.write(f"- {h}")
            else:
                st.write("- No obvious IC-related pages detected – review key sections manually.")
           
# 3) ASSET VERIFICATION (human check of assets & ESG claims)
elif page == "Asset Verification":
    st.header("Asset Verification — Evidence & ESG checks")
    st.caption(
        "This page supports a human verification step for the documents already uploaded. "
        "It helps the Intangibles & Licensing Partner (ILP) or TTO to check ownership, "
        "validity and the robustness of ESG or impact claims before moving to the ILP Console "
        "and reports."
    )

    uploads: List[Any] = ss.get("uploads") or []

    if not uploads:
        st.warning("No documents have been uploaded on the Company page yet.")
    else:
        st.info(
            "Use this view to challenge the evidence: confirm who owns what, whether contracts are in force, "
            "and whether ESG or impact language is backed by real proof (not greenwashing)."
        )

        OPTION_LABELS = ["Select…", "Yes", "No", "Unsure / needs follow-up"]

        for f in uploads:
            fname = getattr(f, "name", "(unnamed file)")
            safe_key = _safe(fname)

            st.markdown(f"#### {fname}")

            col1, col2, col3 = st.columns(3)
            with col1:
                st.selectbox(
                    "Ownership / rights clear?",
                    OPTION_LABELS,
                    key=f"ver_own_{safe_key}",
                    help="Does this document clearly show who owns or controls the asset?",
                )
            with col2:
                st.selectbox(
                    "Up to date & in force?",
                    OPTION_LABELS,
                    key=f"ver_date_{safe_key}",
                    help="Is the document current (still valid, signed, dates make sense)?",
                )
            with col3:
                st.selectbox(
                    "Claims supported by evidence?",
                    OPTION_LABELS,
                    key=f"ver_claim_{safe_key}",
                    help="Where the document makes bold or ESG-related claims, are they backed by real detail?",
                )

            st.text_area(
                "Notes / follow-up for this document",
                key=f"ver_notes_{safe_key}",
                height=80,
                help="Capture anything that needs clarification, extra evidence, or legal review.",
            )

            st.markdown("---")

        st.subheader("Overall verification summary for this company")
        ss["verification_notes"] = st.text_area(
            "Overall verification notes",
            ss.get("verification_notes", ""),
            height=120,
            help="High-level view: which assets look robust, which are weak, and what needs to happen next.",
        )

        st.caption(
            "These checks do **not** replace legal review. They are a structured way for the ILP / TTO to "
            "document how far the evidence can be trusted before licensing design and valuation."
        )
# 4) LIP CONSOLE (was Expert View)
elif page == "LIP Console":
    st.header("LIP Console — Narrative & IC Map")
    nar = st.text_area(
        "Summary (editable by Licensing & Intangibles Partner)",
        value=ss.get("combined_text", ""),
        height=220,
        key="nar_edit",
    )
    ss["narrative"] = nar or ss.get("narrative", "")
    ss["combined_text"] = ss["narrative"]

    colA, colB = st.columns([1, 1])
    with colA:
        if not PUBLIC_MODE:
            st.subheader("Evidence Quality")
            st.progress(min(100, max(0, ss.get("evidence_quality", 0))) / 100.0)
            st.caption(f"{ss.get('evidence_quality', 0)}% evidence coverage (heuristic)")

        st.subheader("4-Leaf Map")
        ic_map: Dict[str, Any] = ss.get("ic_map", {})
        for leaf in ["Human", "Structural", "Customer", "Strategic Alliance"]:
            row = ic_map.get(
                leaf,
                {"tick": False, "narrative": f"No assessment yet for {leaf}.", "score": 0.0},
            )
            tick = "✓" if row["tick"] else "•"
            suffix = "" if PUBLIC_MODE else f"  _(score: {row.get('score', 0.0)})_"
            st.markdown(f"- **{leaf}**: {tick}{suffix}")
            st.caption(row["narrative"])

        st.subheader("Company context (read-only)")
        st.markdown(f"- **Why service:** {ss.get('why_service', '') or '—'}")
        st.markdown(f"- **Stage:** {ss.get('stage', '') or '—'}")
        st.markdown(
            f"- **Plans:** S={ss.get('plan_s', '') or '—'} | "
            f"M={ss.get('plan_m', '') or '—'} | L={ss.get('plan_l', '') or '—'}"
        )
        st.markdown(f"- **Markets & why:** {ss.get('markets_why', '') or '—'}")
        st.markdown(f"- **Target sale & why:** {ss.get('sale_price_why', '') or '—'}")

    with colB:
        st.subheader("Ten-Steps Readiness")
        raw_ten = ss.get("ten_steps") or {}
        scores = raw_ten.get("scores") or [5] * len(TEN_STEPS)
        narrs = raw_ten.get("narratives") or [f"{s}: tbd" for s in TEN_STEPS]
        ten = {"scores": scores, "narratives": narrs}

        st.dataframe(
            {"Step": TEN_STEPS, "Score (1–10)": ten["scores"]},
            hide_index=True,
            use_container_width=True,
        )
        with st.expander("Narrative per step"):
            for s, n in zip(TEN_STEPS, ten["narratives"]):
                st.markdown(f"**{s}** — {n}")
    # ---- VM assumptions sidebar workspace --------------------------------
    try:
        # Sector – adjust key name if you store it differently
        sector = ss.get("sector", ss.get("company_sector", "Unknown"))

        # Build a minimal ic_summary from the ic_map held in session
        ic_map: Dict[str, Any] = ss.get("ic_map", {})
        ic_summary: Dict[str, Dict[str, List[Any]]] = {}

        for cap_label, cap_key in [
            ("Human", "Human"),
            ("Structural", "Structural"),
            ("Customer", "Customer"),
            ("Strategic Alliance", "Strategic"),  # map to Strategic
        ]:
            row = ic_map.get(cap_label, {})
            # crude proxy: 1 explicit item if the tick is True, else none
            explicit_items: List[Any] = ["asset"] if row.get("tick") else []
            ic_summary[cap_key] = {
                "explicit": explicit_items,
                "tacit": [],   # we’re not splitting tacit/explicit here yet
            }

        # Build Ten-Steps scores dict from the stored table
        ten_raw = ss.get("ten_steps", {})
        ten_rows = ten_raw.get("scores") or []
        ten_steps_scores: Dict[str, int] = {
            str(r.get("step")): int(r.get("score", 0)) for r in ten_rows
        }

        vm_assumptions_block(
            sector=sector,
            ic_summary=ic_summary,
            ten_steps_scores=ten_steps_scores,
        )
    except Exception as e:
        st.warning(f"Could not derive assumptions automatically: {e}")
        
# 5) REPORTS
elif page == "Reports":
    st.header("Reports & Exports")
    case_name = ss.get("case_name", "Untitled Company")
    case_folder = OUT_ROOT / _safe(case_name)

def vm_assumptions_block(
    sector: str,
    ic_summary: Dict[str, Dict[str, List[Any]]],
    ten_steps_scores: Dict[str, int],
) -> None:
    """
    Sidebar workspace: VM reviews suggested assumptions, adds any missing ones,
    and confirms what should flow into the IC report.

    Stores:
        st.session_state["vm_assumptions_accepted"]  -> List[VMAssumption]
        st.session_state["vm_assumptions_confirmed"] -> bool
    """
    with st.sidebar:
        st.markdown("### Working assumptions")

        # ---- Suggested assumptions from the engine ------------------------
        suggested = derive_vm_assumptions(sector, ic_summary, ten_steps_scores)

        # Keep any existing manual assumptions in state
        manual_list: List[VMAssumption] = st.session_state.get(
            "vm_manual_assumptions", []
        )

        accepted_suggested: List[VMAssumption] = []

        st.caption("Review and accept the suggested working assumptions:")
        for a in suggested:
            include = st.checkbox(
                a.label,
                value=True,
                key=f"assumption_suggested_{a.key}",
                help=f"Signals: {', '.join(a.source_signals)} | Confidence: {a.confidence}",
            )
            if include:
                a.include = True
                accepted_suggested.append(a)
            else:
                a.include = False

        # ---- Add custom assumptions ---------------------------------------
        st.markdown("---")
        st.caption("Add any missing assumptions:")

        custom_label = st.text_input(
            "Custom assumption title",
            key="custom_assumption_label",
            placeholder="e.g. Early-mover advantage in local market",
        )
        custom_text = st.text_area(
            "Custom assumption narrative",
            key="custom_assumption_text",
            placeholder="Write the assumption in full sentence form.",
        )
        custom_category = st.selectbox(
            "Category",
            ["market", "innovation", "ic-structure", "ten-steps", "other"],
            key="custom_assumption_category",
        )

        if st.button("Add custom assumption"):
            if custom_label.strip() and custom_text.strip():
                new_key = f"manual_{len(manual_list) + 1}"
                manual = VMAssumption(
                    key=new_key,
                    label=custom_label.strip(),
                    narrative=custom_text.strip(),
                    rationale="Added manually by the VM.",
                    category=custom_category,
                    source_signals=["manual_entry"],
                    confidence="high",
                    include=True,
                )
                manual_list.append(manual)
                st.session_state["vm_manual_assumptions"] = manual_list
                st.success("Custom assumption added.")
            else:
                st.warning("Please provide both a title and a narrative.")

        # Show current custom assumptions (read-only list)
        if manual_list:
            st.caption("Custom assumptions added:")
            for m in manual_list:
                st.markdown(f"- ✏️ **{m.label}**")

        # ---- Combine & confirm --------------------------------------------
        # Everything that is either (a) ticked suggested, or (b) custom.
        accepted_all: List[VMAssumption] = accepted_suggested + manual_list
        st.session_state["vm_assumptions_accepted"] = accepted_all

        confirm = st.checkbox(
            "Confirm these assumptions for the IC report",
            key="vm_assumptions_confirmed",
        )

        if confirm and accepted_all:
            st.success("Assumptions confirmed. They will be used in the IC report.")
        elif confirm and not accepted_all:
            st.warning("No assumptions are defined yet. Add or accept at least one.")

    def _compose_ic() -> Tuple[str, str]:
        title = f"IC Report — {case_name}"
        ic_map = ss.get("ic_map", {})

        raw_ten = ss.get("ten_steps") or {}
        scores = raw_ten.get("scores") or [5] * len(TEN_STEPS)
        narrs = raw_ten.get("narratives") or [f"{s}: tbd" for s in TEN_STEPS]
        ten = {"scores": scores, "narratives": narrs}

        b: List[str] = []
        interpreted = ss.get("combined_text", "").strip() or ss.get("narrative", "(no summary)")
        b.append(f"Executive Summary\n\n{interpreted}\n")
        if not PUBLIC_MODE:
            b.append(f"Evidence Quality: ~{ss.get('evidence_quality', 0)}% coverage (heuristic)\n")

        b.append("Four-Leaf Analysis")
        for leaf in ["Human", "Structural", "Customer", "Strategic Alliance"]:
            row = ic_map.get(leaf, {"tick": False, "narrative": "", "score": 0.0})
            tail = "" if PUBLIC_MODE else f" (score: {row.get('score', 0.0)})"
            b.append(f"- {leaf}: {'✓' if row.get('tick') else '•'} — {row.get('narrative', '')}{tail}")

        b.append("\nTen-Steps Readiness")
        for s, n in zip(TEN_STEPS, ten["narratives"]):
            b.append(f"- {n}")

        b.append("\nNotes")
        b.append(
            "This document is provided for high-level evaluation only."
            if PUBLIC_MODE
            else "CONFIDENTIAL. Advisory-first; company and LIP review required for final scoring, licensing design and accounting treatment."
        )
        return title, "\n".join(b)

    def _compose_lic() -> Tuple[str, str]:
        title = f"Licensing Report — {case_name}"
        b: List[str] = []
        b.append(f"Licensing Options & FRAND-Informed Readiness for {case_name}\n")

        b.append("Status & Disclaimer")
        b.append(
            "This report is an advisory draft only. It does not constitute legal advice and must be "
            "reviewed and adapted by qualified legal counsel before signature or implementation.\n"
        )

        b.append("Company Context (selected)")
        b.append(f"- Why service: {ss.get('why_service', '')}")
        b.append(f"- Target sale & why: {ss.get('sale_price_why', '')}\n")

        b.append("Licensing Model Families")
        b.append(
            "- Revenue licences: royalty or fee-based licences (e.g. per unit, per user, revenue share, or per dataset) with "
            "FRAND-informed fee corridors, audit rights and performance conditions."
        )
        b.append(
            "- Access / community licences: royalty-free or low-fee licences that prioritise fair, reasonable "
            "and non-discriminatory access (FRAND-aligned) for social, educational or public-good outcomes."
        )
        b.append(
            "- Co-creation / joint development licences: shared ownership of Foreground IP, clear contribution "
            "records, revenue sharing, and publication rights aligned to partner mandates."
        )
        b.append(
            "- Defensive and cross-licence arrangements: IP pooling, non-assert agreements and mutual access "
            "to codified know-how to reduce litigation risk and accelerate adoption."
        )
        b.append(
            "- Data / algorithm licences: controlled use of datasets, indices, scoring models and algorithms "
            "with clear field-of-use, access tiers and governance.\n"
        )

        b.append("FRAND & Seven Stakeholder Perspective")
        b.append(
            "FRAND is treated here as a design principle rather than a legal certification. For each proposed "
            "licence family, the company should consider fairness, reasonableness and non-discrimination across "
            "the seven stakeholders defined in the Sugai–Weir model (employees, investors, customers, partners "
            "and suppliers, communities and the natural environment). Royalty-bearing and royalty-free licences "
            "can both be FRAND-aligned when access criteria, pricing rationales and governance are clearly stated.\n"
        )

        b.append("Governance & Audit Expectations")
        b.append(
            "- Maintain an IA Register that links each explicit asset (software, indices, datasets, methods, "
            "processes, brand, training content, CRM data) to its licensing model(s)."
        )
        b.append(
            "- Define board-level oversight for licensing, including regular reporting on licence performance, "
            "compliance, ESG and stakeholder impacts."
        )
        b.append(
            "- Ensure that key contracts, JVs, MoUs and access licences are auditable and compatible with "
            "applicable accounting standards (e.g. IAS 38 for intangible assets).\n"
        )

        if PUBLIC_MODE:
            b.append("(Details suppressed in public mode.)")
        return title, "\n".join(b)

    c1, c2 = st.columns(2)
    with c1:
        if st.button("Generate IC Report (DOCX/TXT)", key="btn_ic"):
            title, body = _compose_ic()
            data, fname, mime = _export_bytes(title, body)
            path, msg = _save_bytes(case_folder, fname, data)
            st.download_button(
                "⬇️ Download IC Report",
                data,
                file_name=fname,
                mime=mime,
                key="dl_ic",
            )
            (st.success if path else st.warning)(msg)
    with c2:
        if st.button("Generate Licensing Report (DOCX/TXT)", key="btn_lic"):
            title, body = _compose_lic()
            data, fname, mime = _export_bytes(title, body)
            path, msg = _save_bytes(case_folder, fname, data)
            st.download_button(
                "⬇️ Download Licensing Report",
                data,
                file_name=fname,
                mime=mime,
                key="dl_lic",
            )
            (st.success if path else st.warning)(msg)

    st.caption("Server save root: disabled (public mode)" if PUBLIC_MODE else f"Server save root: {OUT_ROOT}")

# 6) LICENSING TEMPLATES
elif page == "Licensing Templates":
    st.header("Licensing Templates (editable DOCX/TXT)")
    case = ss.get("case_name", "Untitled Company")
    sector = ss.get("sector", "Other")

    template = st.selectbox(
        "Choose a template:",
        ["FRAND Standard", "Co-creation (Joint Development)", "Knowledge (Non-traditional)"],
        index=0,
    )

    if st.button("Generate template", key="btn_make_template"):
        disclaimer = (
            "STATUS: Non-binding draft template. This document is provided for internal discussion only and must be "
            "reviewed and adapted by qualified legal counsel before signature or implementation.\n\n"
        )

        if template == "FRAND Standard":
            title = f"FRAND Standard template — {case}"
            body = (
                f"FRAND Standard Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Purpose & Scope\n"
                "   - Define the licensed technology, dataset, index, software or method.\n"
                "   - Specify fields of use, territories and permitted users.\n\n"
                "2. Access & FRAND-Informed Principles\n"
                "   - Describe how fees (including possible royalty-free tiers) are set on a fair, reasonable and\n"
                "     non-discriminatory basis across comparable users.\n"
                "   - Include access considerations for social, community or public-good partners where relevant.\n\n"
                "3. Financial Terms\n"
                "   - Royalty or fee structure (e.g. per unit, per user, revenue share, or capped fees).\n"
                "   - Invoicing, payment schedule, late payment and currency terms.\n\n"
                "4. Governance, Reporting & Audit\n"
                "   - Licensee reporting obligations (KPIs, usage, sublicensing, ESG/stakeholder indicators if applicable).\n"
                "   - Audit rights and frequency; treatment of under-reporting or non-compliance.\n\n"
                "5. IP Ownership & Improvements\n"
                "   - Background IP ownership and reservation of rights.\n"
                "   - Treatment of improvements, derivative works and feedback.\n\n"
                "6. Compliance, Term & Termination\n"
                "   - Conditions for suspension or termination (breach, non-payment, misuse).\n"
                "   - Survival of key clauses (confidentiality, audit, data protection).\n\n"
                "7. Data, AI & Regulatory Considerations (if applicable)\n"
                "   - Data protection, confidentiality and AI/automated decision-making safeguards.\n"
                "   - Reference to applicable laws, standards or regulatory guidance.\n\n"
                "8. Governing Law & Dispute Resolution\n"
                "   - Governing law (e.g. EU Member State law).\n"
                "   - Mechanism for dispute resolution (negotiation, mediation, arbitration, courts).\n"
            )
        elif template == "Co-creation (Joint Development)":
            title = f"Co-creation template — {case}"
            body = (
                f"Co-creation / Joint Development Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Parties, Purpose & Project Definition\n"
                "   - Identify all parties and describe the joint development project and objectives.\n\n"
                "2. Background IP\n"
                "   - List key Background IP contributed by each party and conditions of use.\n\n"
                "3. Foreground IP & Ownership Structure\n"
                "   - Define Foreground IP and allocate ownership shares (e.g. 50/50 or by contribution).\n"
                "   - Set rules for registration, maintenance and enforcement of Foreground IP.\n\n"
                "4. Contributions, Resources & Cost Sharing\n"
                "   - Describe personnel, facilities, data and funding contributed by each party.\n"
                "   - Agree cost-sharing mechanisms for development and exploitation.\n\n"
                "5. Commercialisation & Revenue Sharing\n"
                "   - Outline commercialisation routes (direct sales, licensing, joint ventures).\n"
                "   - Define revenue sharing mechanisms, including treatment of royalty-free access for selected stakeholders.\n\n"
                "6. Publication, Academic Use & Confidentiality\n"
                "   - Academic and scientific publication rights (timing, review, attribution).\n"
                "   - Confidentiality obligations and carve-outs.\n\n"
                "7. Governance, Decision-Making & Dispute Resolution\n"
                "   - Governance structure (steering committee, decision rules).\n"
                "   - Escalation and dispute resolution framework.\n\n"
                "8. Term, Exit & Transition\n"
                "   - Term and conditions for early termination.\n"
                "   - Exit options (buy-out, assignment, break clauses) and treatment of Foreground IP on exit.\n"
            )
        else:
            title = f"Knowledge licence (non-traditional) — {case}"
            body = (
                f"Knowledge Licence — {case} ({sector})\n\n"
                + disclaimer
                + "1. Knowledge Asset Definition\n"
                "   - Describe the codified know-how (e.g. methods, training materials, playbooks, indices, checklists).\n\n"
                "2. Scope of Licence & Field of Use\n"
                "   - Specify permitted uses (internal training, consulting, product development, policy design, etc.).\n"
                "   - Clarify any restricted uses and prohibited activities.\n\n"
                "3. Access & Stakeholder Considerations\n"
                "   - Define standard and, where appropriate, community or social-benefit access tiers.\n"
                "   - Reference a FRAND-informed approach to fairness, reasonableness and non-discrimination\n"
                "     across employees, investors, customers, partners, suppliers, communities and nature.\n\n"
                "4. Attribution & Moral Rights\n"
                "   - Conditions for attribution (branding, acknowledgements, citation requirements).\n\n"
                "5. Confidentiality, Data Protection & Safeguards\n"
                "   - Treatment of confidential information and personal data.\n"
                "   - Safeguards where AI or automated decision-making is involved.\n\n"
                "6. Term, Revocation & Review\n"
                "   - Duration, renewal and review points.\n"
                "   - Conditions for revocation or modification of the licence.\n\n"
                "7. Governing Law & Dispute Resolution\n"
                "   - Governing law.\n"
                "   - Mechanism for addressing disputes.\n"
            )

        data, fname, mime = _export_bytes(title, body)
        folder = OUT_ROOT / _safe(case)
        path, msg = _save_bytes(folder, fname, data)
        st.download_button(
            "⬇️ Download Template",
            data,
            file_name=fname,
            mime=mime,
            key="dl_tpl",
        )
        (st.success if path else st.warning)(msg)

# 6) LIP ASSISTANT (beta)
elif page == "LIP Assistant":
    st.header("LIP Assistant (beta)")
    st.caption(
        "A lightweight helper for the Licensing & Intangibles Partner. "
        "It re-uses the IC narrative and licensing logic locally — no external AI calls are made in this demo."
    )

    if ss.get("combined_text", ""):
        st.success("IC narrative available from Analyse Evidence / LIP Console.")
    else:
        st.warning("No IC narrative stored yet. Run **Analyse Evidence** first for best results.")

    context_choice = st.selectbox(
        "Which context should the LIP Assistant focus on?",
        ["IC findings", "Licensing options", "Both"],
        index=2,
    )

    question = st.text_area(
        "Your question",
        "",
        height=120,
        help="Example: 'Which assets look IAS 38-ready?' or 'How could we structure a royalty-free community licence?'",
    )

    if st.button("Ask LIP Assistant"):
        if not question.strip():
            st.error("Please enter a question.")
        else:
            ic_text = ss.get("combined_text", "")
            # Short licensing explainer from the report logic
            lic_snippet = (
                "The tool supports revenue licences, access/community licences, co-creation/joint development, "
                "defensive/cross-licence arrangements and data/algorithm licences, all treated through a FRAND-informed "
                "lens across the Seven Stakeholder Model."
            )

            answer_parts: List[str] = []
            answer_parts.append(f"**Question received:** {question.strip()}")

            q_lower = question.lower()

            if context_choice in ("IC findings", "Both"):
                answer_parts.append("**IC & Structural Capital view (IAS 38):**")
                if "structural" in q_lower or "ias 38" in q_lower:
                    answer_parts.append(
                        "- Structural Capital is treated as the home for explicit, documented assets "
                        "(contracts, SOPs, registers, CRM, datasets, board packs). Where both tacit and explicit "
                        "signals exist, Structural is deliberately weighted to dominate, reflecting IAS 38 audit-readiness."
                    )
                if ic_text:
                    answer_parts.append(
                        "- Key IC narrative extract:\n\n"
                        + ic_text[:600]
                        + ("..." if len(ic_text) > 600 else "")
                    )
                else:
                    answer_parts.append(
                        "- No stored IC narrative yet. Run the analysis and revisit this view to see a richer explanation."
                    )

            if context_choice in ("Licensing options", "Both"):
                answer_parts.append("**Licensing & FRAND view:**")
                if "frand" in q_lower:
                    answer_parts.append(
                        "- In this tool, FRAND is a *design principle* rather than a legal certification. "
                        "It guides how pricing, access tiers and governance are framed across different stakeholders, "
                        "including royalty-free and community licences."
                    )
                answer_parts.append("- " + lic_snippet)

            answer_parts.append(
                "_Next step: a Licensing & Intangibles Partner can tailor these insights into concrete clauses using the "
                "reports and templates generated in the other pages._"
            )

            ss["lip_history"].append({"q": question.strip(), "a": "\n\n".join(answer_parts)})
            st.markdown("\n\n".join(answer_parts))

    if ss.get("lip_history"):
        with st.expander("Previous LIP Assistant exchanges"):
            for i, entry in enumerate(reversed(ss["lip_history"]), start=1):
                st.markdown(f"**Q{i}:** {entry['q']}")
                st.markdown(entry["a"])
                st.markdown("---")
